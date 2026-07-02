"""Diagnostic for the .pptx importer (core/pptx_import.parse_pptx).

Builds ONE deck that exercises EVERY styling gap we set out to close, parses it, and prints a
COVERAGE TABLE (feature → CAPTURED / PARTIAL / MISSING + the captured value). python-pptx's
high-level API can't author several features (strike/baseline/hyperlink/buChar/buAutoNum/
gradFill/alpha/outerShdw/headEnd/theme-colored runs/cropped picture/header-filled table) so we
drop to OOXML where needed.

Run:  uv run python -m core.verify_pptx_import
"""

import io
import os
import tempfile

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from core.pptx_import import parse_pptx

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _sub(parent, tag, **attrs):
    el = etree.SubElement(parent, qn(tag))
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


def _png_bytes(color=(200, 80, 80)):
    buf = io.BytesIO()
    Image.new("RGB", (200, 120), color).save(buf, format="PNG")
    return buf.getvalue()


def build_deck(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[5]

    # ---- Slide 1: theme-colored title placeholder (size inheritance) + rich run marks ----
    s1 = prs.slides.add_slide(title_layout)
    title = s1.shapes.title
    title.text = "Titre du thème"
    trun = title.text_frame.paragraphs[0].runs[0]
    rPr = trun._r.get_or_add_rPr()
    # Theme-colored text via schemeClr accent1 + lumMod/lumOff (a real accent shade).
    sf = _sub(rPr, "a:solidFill")
    sc = _sub(sf, "a:schemeClr", val="accent1")
    _sub(sc, "a:lumMod", val="75000")
    _sub(sc, "a:lumOff", val="25000")
    # NOTE: title run has NO explicit sz → must inherit ~44pt from master titleStyle.

    # A second text box exercising every run-level mark.
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "barré "
    r1._r.get_or_add_rPr().set("strike", "sngStrike")
    r2 = p.add_run()
    r2.text = "indice "
    r2._r.get_or_add_rPr().set("baseline", "-25000")
    r3 = p.add_run()
    r3.text = "exposant "
    r3._r.get_or_add_rPr().set("baseline", "30000")
    r4 = p.add_run()
    r4.text = "lien"
    r4.hyperlink.address = "https://example.fr"
    r5 = p.add_run()
    r5.text = " surligné"
    r5rPr = r5._r.get_or_add_rPr()
    hl = _sub(r5rPr, "a:highlight")
    _sub(hl, "a:srgbClr", val="FFFF00")

    # Line spacing on a paragraph (1.5×).
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.5
    rp2 = p2.add_run()
    rp2.text = "interligne 1.5"

    # ---- Slide 2: lists (bullets + numbering + nesting) ----
    s2 = prs.slides.add_slide(blank)
    lb = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    ltf = lb.text_frame
    # Bullet list, two levels.
    bp0 = ltf.paragraphs[0]
    bp0.add_run().text = "Puce niveau 0"
    _bullet(bp0, char=True, level=0)
    _list_para(ltf, "Puce niveau 1", char=True, level=1)
    _list_para(ltf, "Puce niveau 0 again", char=True, level=0)
    # Numbered list (consecutive → one ordered list).
    _list_para(ltf, "Numéro un", char=False, level=0)
    _list_para(ltf, "Numéro deux", char=False, level=0)

    # ---- Slide 3: shapes, gradient, alpha, shadow, rounded corner, stroke / no-stroke ----
    s3 = prs.slides.add_slide(blank)
    # Gradient rectangle.
    grad_shape = s3.shapes.add_shape(
        1, Inches(0.5), Inches(0.5), Inches(3), Inches(2)
    )  # 1=RECTANGLE
    _set_gradient(grad_shape, ["1167D4", "E1000F"], angle=135)
    # Semi-transparent solid + shadow.
    alpha_shape = s3.shapes.add_shape(1, Inches(4), Inches(0.5), Inches(3), Inches(2))
    _set_solid_alpha(alpha_shape, "00AA66", alpha=50000)
    _add_shadow(alpha_shape)
    # Rounded rectangle with explicit adjustment.
    round_shape = s3.shapes.add_shape(
        5, Inches(8), Inches(0.5), Inches(3), Inches(2)
    )  # 5=ROUNDED_RECTANGLE
    _set_round_adj(round_shape, 20000)
    # A plain shape with NO explicit outline → must NOT get a spurious stroke (gap #5a).
    no_border = s3.shapes.add_shape(1, Inches(0.5), Inches(3.5), Inches(3), Inches(2))
    _set_solid(no_border, "888888")
    _strip_line(no_border)

    # ---- Slide 4: connector with arrowheads + cropped picture + styled table ----
    s4 = prs.slides.add_slide(blank)
    conn = s4.shapes.add_connector(
        2, Inches(0.5), Inches(0.5), Inches(5), Inches(0.5)
    )  # 2=STRAIGHT
    _set_arrowheads(conn)
    # Cropped picture.
    pic = s4.shapes.add_picture(
        io.BytesIO(_png_bytes()), Inches(0.5), Inches(1.5), Inches(3), Inches(2)
    )
    pic.crop_left = 0.1
    pic.crop_right = 0.05
    pic.crop_top = 0.2
    pic.crop_bottom = 0.0
    # Styled table: header-row fill + border + banding.
    tbl_shape = s4.shapes.add_table(3, 2, Inches(5), Inches(1.5), Inches(6), Inches(3))
    table = tbl_shape.table
    table.cell(0, 0).text = "En-tête A"
    table.cell(0, 1).text = "En-tête B"
    table.cell(1, 0).text = "x"
    _style_table(tbl_shape)

    # ---- Slide 5: empty placeholder (must be SKIPPED) + theme-filled shape ----
    s5 = prs.slides.add_slide(title_layout)
    # leave the title empty → empty text object, should be skipped.
    themed = s5.shapes.add_shape(1, Inches(1), Inches(3), Inches(4), Inches(2))
    _set_scheme_fill(themed, "accent2")

    prs.save(path)


# --- OOXML authoring helpers -------------------------------------------------


def _bullet(para, char, level):
    para.level = level
    pPr = para._p.get_or_add_pPr()
    if char:
        _sub(pPr, "a:buChar", char="•")
    else:
        _sub(pPr, "a:buAutoNum", type="arabicPeriod")


def _list_para(tf, text, char, level):
    """Append a new paragraph with `text` as a bullet (char) or numbered (auto) list item."""
    para = tf.add_paragraph()
    para.add_run().text = text
    _bullet(para, char=char, level=level)
    return para


def _spPr(shape):
    return shape._element.spPr


def _strip_existing_fill(spPr):
    for tag in (
        "a:noFill",
        "a:solidFill",
        "a:gradFill",
        "a:blipFill",
        "a:pattFill",
        "a:grpFill",
    ):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)


def _set_solid(shape, hex6):
    spPr = _spPr(shape)
    _strip_existing_fill(spPr)
    sf = _sub(spPr, "a:solidFill")
    _sub(sf, "a:srgbClr", val=hex6)


def _set_solid_alpha(shape, hex6, alpha):
    spPr = _spPr(shape)
    _strip_existing_fill(spPr)
    sf = _sub(spPr, "a:solidFill")
    c = _sub(sf, "a:srgbClr", val=hex6)
    _sub(c, "a:alpha", val=alpha)


def _set_scheme_fill(shape, scheme_val):
    spPr = _spPr(shape)
    _strip_existing_fill(spPr)
    sf = _sub(spPr, "a:solidFill")
    _sub(sf, "a:schemeClr", val=scheme_val)


def _set_gradient(shape, hexes, angle):
    spPr = _spPr(shape)
    _strip_existing_fill(spPr)
    grad = _sub(spPr, "a:gradFill")
    gsLst = _sub(grad, "a:gsLst")
    n = len(hexes)
    for i, h in enumerate(hexes):
        pos = round(i / (n - 1) * 100000)
        gs = _sub(gsLst, "a:gs", pos=pos)
        _sub(gs, "a:srgbClr", val=h)
    # OOXML angle is 1/60000 degree, clockwise from 3 o'clock.
    _sub(grad, "a:lin", ang=round((angle - 90) % 360 * 60000), scaled="1")


def _add_shadow(shape):
    spPr = _spPr(shape)
    eff = _sub(spPr, "a:effectLst")
    sh = _sub(eff, "a:outerShdw", blurRad="40000", dist="20000", dir="5400000")
    _sub(sh, "a:srgbClr", val="000000")


def _set_round_adj(shape, val):
    spPr = _spPr(shape)
    prstGeom = spPr.find(qn("a:prstGeom"))
    avLst = prstGeom.find(qn("a:avLst"))
    if avLst is None:
        avLst = _sub(prstGeom, "a:avLst")
    gd = _sub(avLst, "a:gd")
    gd.set("name", "adj")
    gd.set("fmla", f"val {val}")


def _strip_line(shape):
    spPr = _spPr(shape)
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        spPr.remove(ln)
    # explicit noFill line to be sure nothing inherits.
    ln = _sub(spPr, "a:ln")
    _sub(ln, "a:noFill")


def _set_arrowheads(conn):
    spPr = _spPr(conn)
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = _sub(spPr, "a:ln")
    ln.set("w", "38100")
    sf = _sub(ln, "a:solidFill")
    _sub(sf, "a:srgbClr", val="333333")
    _sub(ln, "a:headEnd", type="triangle")
    _sub(ln, "a:tailEnd", type="stealth")


def _style_table(tbl_shape):
    tbl = tbl_shape._element.find(".//" + qn("a:tbl"))
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = _sub(tbl, "a:tblPr")
    tblPr.set("bandRow", "1")
    tblPr.set("firstRow", "1")
    # Header-row fill on first row's cells + a border on the first cell.
    first_tr = tbl.find(qn("a:tr"))
    for i, tc in enumerate(first_tr.findall(qn("a:tc"))):
        tcPr = tc.find(qn("a:tcPr"))
        if tcPr is None:
            tcPr = _sub(tc, "a:tcPr")
        if i == 0:
            lnL = _sub(tcPr, "a:lnL", w="12700")
            sfl = _sub(lnL, "a:solidFill")
            _sub(sfl, "a:srgbClr", val="2F5496")
        sf = _sub(tcPr, "a:solidFill")
        _sub(sf, "a:srgbClr", val="4472C4")


# --- coverage reporting ------------------------------------------------------


def _find_text(slides, contains):
    for s in slides:
        for o in s["objects"]:
            if o.get("type") == "text":
                for para in o.get("paragraphs", []):
                    for run in para.get("runs", []):
                        if contains in (run.get("text") or ""):
                            return o, para, run
    return None, None, None


def _find(slides, pred):
    for s in slides:
        for o in s["objects"]:
            if pred(o):
                return o
    return None


def report(result):
    slides = result["slides"]
    rows = []  # (tier, feature, status, value)

    def add(tier, feature, status, value=""):
        rows.append((tier, feature, status, str(value)))

    # Tier 1
    _, _, r = _find_text(slides, "barré")
    add(
        1,
        "strikethrough",
        "CAPTURED" if r and r.get("strike") else "MISSING",
        r and r.get("strike"),
    )
    _, _, sub = _find_text(slides, "indice")
    _, _, sup = _find_text(slides, "exposant")
    sub_ok = sub and sub.get("sub")
    sup_ok = sup and sup.get("sup")
    add(
        1,
        "subscript / superscript",
        "CAPTURED" if sub_ok and sup_ok else "MISSING",
        f"sub={sub and sub.get('sub')} sup={sup and sup.get('sup')}",
    )
    _, _, lk = _find_text(slides, "lien")
    add(
        1,
        "hyperlink",
        "CAPTURED" if lk and lk.get("link") else "MISSING",
        lk and lk.get("link"),
    )
    _, _, hlr = _find_text(slides, "surligné")
    add(
        1,
        "highlight",
        "CAPTURED" if hlr and hlr.get("highlight") else "MISSING",
        hlr and hlr.get("highlight"),
    )
    # 5a: no spurious stroke. The no-border gray shape must have no stroke/strokeWidth.
    gray = _find(
        slides, lambda o: o.get("type") == "shape" and o.get("fill") == "#888888"
    )
    no_spurious = (
        gray is not None and "stroke" not in gray and "strokeWidth" not in gray
    )
    add(
        1,
        "cleanup: no spurious border",
        "CAPTURED" if no_spurious else "MISSING",
        f"stroke={gray and gray.get('stroke')} w={gray and gray.get('strokeWidth')}",
    )
    # 5b: empty placeholder skipped → slide 5 should have NO text object.
    s5_texts = [o for o in slides[4]["objects"] if o.get("type") == "text"]
    add(
        1,
        "cleanup: skip empty text",
        "CAPTURED" if not s5_texts else "MISSING",
        f"empty text objs on slide5={len(s5_texts)}; skipped={result['skipped'].get('empty_text')}",
    )

    # Tier 2
    ls_obj, _, _ = _find_text(slides, "interligne")
    add(
        2,
        "line spacing → lineHeight",
        "CAPTURED" if ls_obj and ls_obj.get("lineHeight") else "MISSING",
        ls_obj and ls_obj.get("lineHeight"),
    )
    img = _find(slides, lambda o: o.get("type") == "image")
    crop = img and img.get("crop")
    add(2, "image crop (percent)", "CAPTURED" if crop else "MISSING", crop)
    grad = _find(slides, lambda o: o.get("type") == "shape" and o.get("gradient"))
    add(
        2,
        "shape gradient fill",
        "CAPTURED" if grad else "MISSING",
        grad and grad.get("gradient"),
    )
    alpha = _find(
        slides, lambda o: o.get("type") == "shape" and o.get("opacity") is not None
    )
    add(
        2,
        "shape opacity (alpha)",
        "CAPTURED" if alpha else "MISSING",
        alpha and alpha.get("opacity"),
    )
    shadow = _find(slides, lambda o: o.get("type") == "shape" and o.get("shadow"))
    add(
        2,
        "shape shadow",
        "CAPTURED" if shadow else "MISSING",
        shadow and shadow.get("shadow"),
    )
    rrect = _find(slides, lambda o: o.get("shape") == "roundRect" and o.get("radius"))
    add(
        2,
        "shape corner radius",
        "CAPTURED" if rrect else "MISSING",
        rrect and rrect.get("radius"),
    )
    line = _find(slides, lambda o: o.get("type") == "line")
    arrows_ok = line and line.get("arrowStart") and line.get("arrowEnd")
    add(
        2,
        "line arrowheads",
        "CAPTURED" if arrows_ok else "MISSING",
        f"start={line and line.get('arrowStart')} end={line and line.get('arrowEnd')}",
    )
    table = _find(slides, lambda o: o.get("type") == "table")
    tbl_ok = table and (
        table.get("fill") or table.get("stroke") or table.get("banding")
    )
    tbl_full = (
        table and table.get("fill") and table.get("stroke") and table.get("banding")
    )
    add(
        2,
        "table styling (fill/stroke/banding)",
        "CAPTURED" if tbl_full else ("PARTIAL" if tbl_ok else "MISSING"),
        table and {k: table.get(k) for k in ("fill", "stroke", "banding")},
    )

    # Tier 3
    bullet = _find(
        slides,
        lambda o: o.get("type") == "text"
        and any(p.get("list") == "bullet" for p in o.get("paragraphs", [])),
    )
    number = _find(
        slides,
        lambda o: o.get("type") == "text"
        and any(p.get("list") == "number" for p in o.get("paragraphs", [])),
    )
    nested = bullet and any(
        p.get("list") == "bullet" and p.get("level") == 1
        for p in bullet.get("paragraphs", [])
    )
    lists_status = (
        "CAPTURED"
        if (bullet and number and nested)
        else ("PARTIAL" if (bullet or number) else "MISSING")
    )
    add(
        3,
        "lists (bullet/number/indent)",
        lists_status,
        f"bullet={bool(bullet)} number={bool(number)} nestedLvl1={bool(nested)}",
    )
    # theme color: title run color resolved from accent1+lumMod/lumOff.
    title_obj, _, title_run = _find_text(slides, "Titre du thème")
    theme_color = title_run and title_run.get("color")
    themed_shape = _find(
        slides,
        lambda o: o.get("type") == "shape" and o.get("fill") and o.get("w") == 384,
    )
    add(
        3,
        "theme colors (run + shape)",
        "CAPTURED" if theme_color else "MISSING",
        f"runColor={theme_color} shapeFill(accent2)={themed_shape and themed_shape.get('fill')}",
    )
    # size inheritance: title run/box should be near 44pt → ~59px (44 * 96/72).
    title_size = (title_run and title_run.get("fontSize")) or (
        title_obj and title_obj.get("fontSize")
    )
    size_ok = title_size and title_size >= 40  # well above the ~22px editor default
    add(
        3,
        "inherited placeholder font size",
        "CAPTURED" if size_ok else ("PARTIAL" if title_size else "MISSING"),
        f"title fontSize px={title_size} (expect ~59)",
    )

    # ---- print table ----
    print("\n" + "=" * 78)
    print("PPTX IMPORT COVERAGE")
    print("=" * 78)
    cur_tier = None
    n_missing = 0
    for tier, feature, status, value in rows:
        if tier != cur_tier:
            print(f"\n--- Tier {tier} ---")
            cur_tier = tier
        if status == "MISSING":
            n_missing += 1
        val = value if len(value) <= 46 else value[:43] + "..."
        print(f"  {status:9}  {feature:34}  {val}")
    print("\n" + "-" * 78)
    print(f"skipped (diagnostic): {result.get('skipped')}")
    print(f"slideSize: {result['slideSize']}  slides: {len(slides)}")
    print("-" * 78)
    if n_missing:
        print(f"RESULT: {n_missing} feature(s) MISSING")
    else:
        print("RESULT: every feature CAPTURED or PARTIAL (none MISSING)")
    return n_missing


def main():
    with tempfile.TemporaryDirectory() as tmp:
        path = os.path.join(tmp, "diagnostic.pptx")
        build_deck(path)
        result = parse_pptx(path)
    n_missing = report(result)
    raise SystemExit(1 if n_missing else 0)


if __name__ == "__main__":
    main()
