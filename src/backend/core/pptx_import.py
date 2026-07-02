"""Real .pptx → native scene-graph parser (MIT `python-pptx`, fully local, no cloud).

Turns a PowerPoint file into the structured object views our frontend understands (text boxes
with rich runs, shapes, pictures, tables, lines) so import produces EDITABLE objects rather than
a flat PNG per slide. The response contract (consumed by ImportDialog.tsx):

    { "mode": "objects",
      "slideSize": { "w": <px>, "h": <px> },
      "slides": [ { "background": "#rrggbb"|null, "objects": [ <object view> ] }, ... ] }

Units: PowerPoint measures in EMU (914400 EMU/inch). We render at 96 px/inch, so
**px = EMU / 9525**. Font sizes are in points (Pt); CSS px = pt * 96 / 72. Colors are RGB hex.

Styling beyond what python-pptx's high-level API exposes (theme colors, gradients, alpha,
shadows, bullets, baseline/strike/hyperlink/highlight) is read straight off the underlying
OOXML (`rPr`/`spPr`/`pPr`) via lxml. A per-presentation `ThemeResolver` turns scheme colors
(accent1, tx1, …), honoring the master `clrMap` and lumMod/lumOff/tint/shade, into hex.

Anything we cannot map (unknown auto-shape, group, OLE, chart, …) is skipped gracefully and
counted, never raised, the caller falls back to the LibreOffice→PNG path on a hard failure.
"""

import base64
import colorsys
import math

import lxml.etree as etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

# Hardened parser for XML taken straight from the uploaded archive (theme parts):
# entity resolution and network access stay off so a crafted deck can't XXE us.
XML_PARSER = etree.XMLParser(resolve_entities=False, no_network=True)

# EMU per CSS px at 96 dpi (914400 EMU/in ÷ 96 px/in).
EMU_PER_PX = 9525
# Pt → px at 96 dpi (96/72).
PT_TO_PX = 96.0 / 72.0

# PowerPoint's default text-frame insets (bodyPr) when lIns/tIns/rIns/bIns are unset:
# left/right = 0.1in = 91440 EMU, top/bottom = 0.05in = 45720 EMU. We convert to px so the
# text box padding matches PowerPoint (and the LibreOffice ground truth) instead of a guess.
_DEFAULT_INS = {"lIns": 91440, "tIns": 45720, "rIns": 91440, "bIns": 45720}

# PowerPoint's "single" line spacing is the font's natural line box (~1.2× font size for the body
# fonts here), not CSS's looser default (~1.5). A:lnSpc spcPct is a multiple of THIS single spacing,
# so a 140% spcPct ≈ 1.4 × 1.2 in CSS line-height terms. We apply this factor to both the default and
# percent line spacing so multi-line text occupies the same vertical extent as the ground truth.
_SINGLE_LINE_FACTOR = 1.1


def emu_to_px(value):
    """EMU → integer px (rounded). None-safe (returns 0)."""
    if value is None:
        return 0
    return round(int(value) / EMU_PER_PX)


def _hex(color):
    """A python-pptx RGBColor → "#rrggbb", or None if not a usable solid RGB color.

    Theme colors (MSO_THEME_COLOR) have no resolvable RGB without the theme part. We resolve
    those separately from the XML via ThemeResolver; this helper only handles explicit RGB.
    """
    if color is None:
        return None
    try:
        # .type is None for an unset color; THEME colors raise on .rgb.
        if color.type is None:
            return None
        rgb = color.rgb  # raises for theme-only colors
    except Exception:
        return None
    if rgb is None:
        return None
    return "#" + str(rgb).lower()


# ---------------------------------------------------------------------------
# Theme color resolution (gap #15)
# ---------------------------------------------------------------------------

# Sentinels for the two system colors PowerPoint themes use for dk1/lt1.
_SYS_DEFAULTS = {"windowText": "000000", "window": "ffffff"}

# MSO_THEME_COLOR enum name → the clrMap slot it refers to. The clrMap (read off the master)
# then redirects that slot to a real scheme entry (e.g. tx1 → dk1, bg2 → lt2).
_THEME_NAME_TO_SLOT = {
    "DARK_1": "dk1",
    "LIGHT_1": "lt1",
    "DARK_2": "dk2",
    "LIGHT_2": "lt2",
    "TEXT_1": "tx1",
    "TEXT_2": "tx2",
    "BACKGROUND_1": "bg1",
    "BACKGROUND_2": "bg2",
    "ACCENT_1": "accent1",
    "ACCENT_2": "accent2",
    "ACCENT_3": "accent3",
    "ACCENT_4": "accent4",
    "ACCENT_5": "accent5",
    "ACCENT_6": "accent6",
    "HYPERLINK": "hlink",
    "FOLLOWED_HYPERLINK": "folHlink",
}


def _clamp(v):
    return max(0.0, min(1.0, v))


def _hex_to_rgb(h):
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _rgb_to_hex(r, g, b):
    return "#%02x%02x%02x" % (
        round(_clamp(r / 255) * 255),
        round(_clamp(g / 255) * 255),
        round(_clamp(b / 255) * 255),
    )


def _apply_lum_adjust(hex6, adjustments):
    """Apply lumMod/lumOff/tint/shade (each a fraction 0..1) to a base "#rrggbb".

    We approximate PowerPoint's HSL-luminance model: lumMod scales luminance, lumOff adds to it,
    tint blends toward white, shade blends toward black. Good enough that accent shades land in
    the right ballpark (the exact transform is in scRGB space; this is the common approximation).
    """
    r, g, b = _hex_to_rgb(hex6)
    rf, gf, bf = r / 255.0, g / 255.0, b / 255.0

    lum_mod = adjustments.get("lumMod")
    lum_off = adjustments.get("lumOff")
    tint = adjustments.get("tint")
    shade = adjustments.get("shade")

    if lum_mod is not None or lum_off is not None:
        hue, lum, sat = colorsys.rgb_to_hls(rf, gf, bf)
        if lum_mod is not None:
            lum = lum * lum_mod
        if lum_off is not None:
            lum = lum + lum_off
        lum = _clamp(lum)
        rf, gf, bf = colorsys.hls_to_rgb(hue, lum, sat)

    if tint is not None:
        # tint: blend toward white by (1 - tint).
        rf = rf * tint + (1.0 - tint)
        gf = gf * tint + (1.0 - tint)
        bf = bf * tint + (1.0 - tint)
    if shade is not None:
        # shade: scale toward black.
        rf, gf, bf = rf * shade, gf * shade, bf * shade

    return _rgb_to_hex(round(rf * 255), round(gf * 255), round(bf * 255))


class ThemeResolver:
    """Resolves theme/scheme colors to hex for one presentation, honoring the master clrMap.

    Built lazily per slide from the slide's master (decks can have several masters/themes).
    Caches the resolved {slot → hex} scheme so repeated lookups are cheap.
    """

    def __init__(self, scheme, clrmap, major_latin=None, minor_latin=None):
        self._scheme = (
            scheme  # {'dk1':'000000', 'accent1':'4f81bd', ...} (no '#', lowercase)
        )
        self._clrmap = clrmap  # {'tx1':'dk1', 'bg1':'lt1', ...}
        # Theme major (headings) / minor (body) latin typeface, what +mj-lt / +mn-lt resolve to.
        # Runs that carry no explicit <a:latin> inherit these (python-pptx's font.name returns None),
        # so resolving them here lets theme-Calibri text map to Carlito and match the ground truth.
        self._major_latin = major_latin
        self._minor_latin = minor_latin

    @classmethod
    def for_slide(cls, slide, cache):
        """Get (memoized) the resolver for a slide's master+theme. `cache` is a dict keyed by id."""
        try:
            master = slide.slide_layout.slide_master
        except Exception:
            return cls({}, {})
        key = id(master)
        if key in cache:
            return cache[key]
        resolver = cls._build(master)
        cache[key] = resolver
        return resolver

    @classmethod
    def _build(cls, master):
        scheme = {}
        clrmap = {}
        major_latin = None
        minor_latin = None
        try:
            theme_part = master.part.part_related_by(RT.THEME)

            root = etree.fromstring(theme_part.blob, parser=XML_PARSER)
            clr_scheme = root.find(".//" + qn("a:clrScheme"))
            if clr_scheme is not None:
                for entry in clr_scheme:
                    slot = etree.QName(entry).localname  # dk1, lt1, accent1, ...
                    hex6 = cls._read_color_child(entry)
                    if hex6:
                        scheme[slot] = hex6
            # Theme font scheme: the latin typeface for major (headings) and minor (body) fonts.
            for tag, attr in (("a:majorFont", "major"), ("a:minorFont", "minor")):
                font_el = root.find(".//" + qn(tag))
                if font_el is None:
                    continue
                latin = font_el.find(qn("a:latin"))
                tf = latin.get("typeface") if latin is not None else None
                if tf:
                    if attr == "major":
                        major_latin = tf
                    else:
                        minor_latin = tf
        except Exception:
            pass
        try:
            cm = master.element.find(qn("p:clrMap"))
            if cm is not None:
                clrmap = {k: v for k, v in cm.attrib.items()}
        except Exception:
            pass
        return cls(scheme, clrmap, major_latin, minor_latin)

    def theme_font(self, role):
        """The theme latin typeface a run inherits when it carries no explicit font.

        role 'title' -> major (heading) font; everything else -> minor (body) font. Returns the
        typeface name (e.g. 'Calibri') or None. PowerPoint may store these as '+mj-lt'/'+mn-lt'
        references too; we already dereferenced them to the literal theme typeface, but guard anyway.
        """
        name = self._major_latin if role == "title" else self._minor_latin
        if name and name.startswith("+"):
            # An unexpected unresolved reference, map the well-known tokens to our captured fonts.
            name = self._major_latin if "mj" in name else self._minor_latin
        return name

    @staticmethod
    def _read_color_child(entry):
        """First color child of a clrScheme entry → lowercase 6-hex (handles srgbClr/sysClr)."""
        for child in entry:
            tag = etree.QName(child).localname
            if tag == "srgbClr":
                val = child.get("val")
                if val:
                    return val.lower()
            elif tag == "sysClr":
                last = child.get("lastClr")
                if last:
                    return last.lower()
                return _SYS_DEFAULTS.get(child.get("val"), "000000")
        return None

    def scheme_hex(self, slot):
        """A clrMap slot or scheme name (e.g. 'tx1', 'accent1') → "#rrggbb" or None."""
        resolved = self._clrmap.get(slot, slot)
        hex6 = self._scheme.get(resolved)
        return "#" + hex6 if hex6 else None

    def resolve_theme_color(self, mso_theme_color, adjustments=None):
        """An MSO_THEME_COLOR (enum) → "#rrggbb" with optional lum/tint/shade applied."""
        name = getattr(mso_theme_color, "name", None)
        slot = _THEME_NAME_TO_SLOT.get(name)
        if not slot:
            return None
        base = self.scheme_hex(slot)
        if not base:
            return None
        if adjustments:
            return _apply_lum_adjust(base[1:], adjustments)
        return base

    def resolve_scheme_element(self, scheme_clr_el):
        """An <a:schemeClr> lxml element → "#rrggbb", honoring its lum/tint/shade children."""
        val = scheme_clr_el.get(
            "val"
        )  # e.g. 'accent1', 'tx1', 'phClr' (placeholder, skip)
        if not val or val == "phClr":
            return None
        base = self.scheme_hex(val)
        if not base:
            return None
        adj = _read_color_adjustments(scheme_clr_el)
        if adj:
            return _apply_lum_adjust(base[1:], adj)
        return base


def _read_color_adjustments(color_el):
    """Read lumMod/lumOff/tint/shade children off a color element → {name: fraction 0..1}."""
    out = {}

    for child in color_el:
        tag = etree.QName(child).localname
        if tag in ("lumMod", "lumOff", "tint", "shade"):
            val = child.get("val")
            if val is not None:
                try:
                    out[tag] = int(val) / 100000.0
                except ValueError:
                    pass
    return out


def _resolve_color_element(color_el, theme):
    """Any color element (<a:srgbClr> or <a:schemeClr>) → "#rrggbb" or None.

    srgbClr carries an explicit hex (with optional adjustments); schemeClr goes through the theme.
    """
    if color_el is None:
        return None

    tag = etree.QName(color_el).localname
    if tag == "srgbClr":
        val = color_el.get("val")
        if not val:
            return None
        adj = _read_color_adjustments(color_el)
        return _apply_lum_adjust(val.lower(), adj) if adj else "#" + val.lower()
    if tag == "schemeClr" and theme is not None:
        return theme.resolve_scheme_element(color_el)
    return None


def _first_color_el(parent):
    """First child of `parent` that is a color element (srgbClr/schemeClr/sysClr), or None."""
    if parent is None:
        return None

    for child in parent:
        tag = etree.QName(child).localname
        if tag in ("srgbClr", "schemeClr", "sysClr"):
            return child
    return None


# PowerPoint paragraph alignment → our 'left'|'center'|'right'|'justify'.
_ALIGN = {
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "justify",
    PP_ALIGN.JUSTIFY_LOW: "justify",
    PP_ALIGN.DISTRIBUTE: "justify",
}

# auto_shape_type (MSO_AUTO_SHAPE_TYPE) name → our ShapeKind. We match on the enum *name* so we
# don't need to import the (large) enum; unmapped shapes fall back to 'rect'.
_SHAPE_KIND = {
    "RECTANGLE": "rect",
    "ROUNDED_RECTANGLE": "roundRect",
    "ROUND_1_RECTANGLE": "roundRect",
    "ROUND_2_SAME_RECTANGLE": "roundRect",
    "OVAL": "ellipse",
    "ISOCELES_TRIANGLE": "triangle",
    "RIGHT_TRIANGLE": "triangle",
    "DIAMOND": "diamond",
    "PENTAGON": "pentagon",
    "REGULAR_PENTAGON": "pentagon",
    "HEXAGON": "hexagon",
    "STAR_5_POINT": "star",
    "STAR_5_POINTS": "star",
    "RIGHT_ARROW": "arrowRight",
    "LEFT_ARROW": "arrowLeft",
    "CHEVRON": "chevron",
    "RECTANGULAR_CALLOUT": "callout",
    "ROUNDED_RECTANGULAR_CALLOUT": "callout",
    "OVAL_CALLOUT": "callout",
}

# roundRect shape kinds for which a corner radius is meaningful.
_ROUND_KINDS = {"roundRect"}


def _vanchor(text_frame):
    """Vertical anchor of a text frame → our valign ('top'|'middle'|'bottom')."""
    try:
        anchor = text_frame.vertical_anchor
    except Exception:
        anchor = None
    if anchor is None:
        return "top"
    name = getattr(anchor, "name", str(anchor)) or ""
    if "BOTTOM" in name:
        return "bottom"
    if "MIDDLE" in name or "CENTER" in name:
        return "middle"
    return "top"


def _body_pr(shape):
    """The shape's <a:bodyPr> element (text-frame properties), or None."""
    try:
        el = getattr(shape, "_element", None)
        txBody = el.find(qn("p:txBody")) if el is not None else None
        return txBody.find(qn("a:bodyPr")) if txBody is not None else None
    except Exception:
        return None


def _text_insets(body_pr):
    """bodyPr lIns/tIns/rIns/bIns (EMU) → padding px {top,right,bottom,left}.

    Uses PowerPoint's defaults (l/r=0.1in, t/b=0.05in) for any inset the author left unset, so the
    text starts where PowerPoint puts it, and the wrap width matches, instead of our old fixed
    8px/12px guess. A bodyPr that explicitly sets lIns="0" therefore yields 0 px padding.
    """
    out = {}
    for emu_key, css_key in (
        ("tIns", "top"),
        ("rIns", "right"),
        ("bIns", "bottom"),
        ("lIns", "left"),
    ):
        emu = _DEFAULT_INS[emu_key]
        if body_pr is not None:
            v = body_pr.get(emu_key)
            if v is not None:
                try:
                    emu = int(v)
                except ValueError:
                    pass
        out[css_key] = emu_to_px(emu)
    return out


def _autofit_kind(body_pr):
    """bodyPr autofit child → 'norm' (shrink text), 'shape' (grow box), or None (no autofit)."""
    if body_pr is None:
        return None
    if body_pr.find(qn("a:normAutofit")) is not None:
        return "norm"
    if body_pr.find(qn("a:spAutoFit")) is not None:
        return "shape"
    return None


def _norm_autofit_scale(body_pr):
    """normAutofit fontScale/lnSpcReduction → (font_scale 0..1, ln_spc_reduction 0..1).

    PowerPoint stores these as 1/1000 of a percent (e.g. fontScale="92500" = 92.5%). It bakes the
    shrink into the layout, so to match the ground truth we must scale every run's effective size by
    font_scale and reduce line spacing by ln_spc_reduction. Returns (1.0, 0.0) when unset.
    """
    font_scale = 1.0
    ln_reduction = 0.0
    if body_pr is None:
        return font_scale, ln_reduction
    na = body_pr.find(qn("a:normAutofit"))
    if na is None:
        return font_scale, ln_reduction
    fs = na.get("fontScale")
    if fs is not None:
        try:
            font_scale = int(fs) / 100000.0
        except ValueError:
            pass
    lr = na.get("lnSpcReduction")
    if lr is not None:
        try:
            ln_reduction = int(lr) / 100000.0
        except ValueError:
            pass
    return font_scale, ln_reduction


def _para_spacing(para):
    """Per-paragraph spacing off a:pPr → {lineHeight?, lineHeightPx?, spaceBefore?, spaceAfter?}.

    - a:lnSpc/a:spcPct (percent, /1000%)  → lineHeight unitless multiple (140000 → 1.4)
    - a:lnSpc/a:spcPts (points, /100pt)   → lineHeightPx (absolute px line box)
    - a:spcBef/a:spcPts                   → spaceBefore px (margin-top)
    - a:spcAft/a:spcPts                   → spaceAfter px  (margin-bottom)
    Empty when the paragraph sets none. Only spcPts is meaningful for spcBef/spcAft in practice.
    """
    out = {}
    try:
        pPr = para._p.find(qn("a:pPr"))
    except Exception:
        pPr = None
    if pPr is None:
        return out

    ln = pPr.find(qn("a:lnSpc"))
    if ln is not None:
        pct = ln.find(qn("a:spcPct"))
        pts = ln.find(qn("a:spcPts"))
        if pct is not None and pct.get("val"):
            try:
                # spcPct is a multiple of single spacing → scale by the natural single-line factor
                # so 140% lands at ~1.68 CSS line-height, matching PowerPoint's actual line box.
                out["lineHeight"] = round(
                    int(pct.get("val")) / 100000.0 * _SINGLE_LINE_FACTOR, 4
                )
            except ValueError:
                pass
        elif pts is not None and pts.get("val"):
            try:
                out["lineHeightPx"] = round(int(pts.get("val")) / 100.0 * PT_TO_PX, 2)
            except ValueError:
                pass

    for tag, key in (("a:spcBef", "spaceBefore"), ("a:spcAft", "spaceAfter")):
        el = pPr.find(qn(tag))
        if el is None:
            continue
        pts = el.find(qn("a:spcPts"))
        if pts is not None and pts.get("val"):
            try:
                out[key] = round(int(pts.get("val")) / 100.0 * PT_TO_PX, 2)
            except ValueError:
                pass
    return out


def _para_indent(para):
    """Per-paragraph marL/indent off a:pPr (EMU) → {marginLeft px, hangingIndent px}.

    marL is the left margin of the whole paragraph; indent is the first-line delta (negative for a
    hanging bullet). We surface marL so nested levels indent like the source even when the marL is
    authored explicitly rather than implied by the level.
    """
    out = {}
    try:
        pPr = para._p.find(qn("a:pPr"))
    except Exception:
        pPr = None
    if pPr is None:
        return out
    marL = pPr.get("marL")
    if marL is not None:
        try:
            out["marginLeft"] = emu_to_px(int(marL))
        except ValueError:
            pass
    indent = pPr.get("indent")
    if indent is not None:
        try:
            out["indent"] = emu_to_px(int(indent))
        except ValueError:
            pass
    return out


def _run_color(run, theme):
    """Resolve a run's text color to "#rrggbb", explicit RGB or theme-mapped, or None.

    Reads straight off the run's rPr/solidFill so theme colors (schemeClr) and adjustments
    (lumMod/lumOff/tint/shade) resolve; falls back to python-pptx's RGB for plain srgbClr.
    """
    try:
        rPr = run._r.find(qn("a:rPr"))
        if rPr is not None:
            fill = rPr.find(qn("a:solidFill"))
            if fill is not None:
                color_el = _first_color_el(fill)
                resolved = _resolve_color_element(color_el, theme)
                if resolved:
                    return resolved
    except Exception:
        pass
    # Fallback: explicit RGB via the high-level API (covers cases without our XML path).
    return _hex(getattr(run.font, "color", None))


def _run_marks(run):
    """Read run-level marks off rPr XML: strike, sub/sup (baseline), highlight. Returns a dict."""
    out = {}
    try:
        rPr = run._r.find(qn("a:rPr"))
    except Exception:
        rPr = None
    if rPr is None:
        return out
    strike = rPr.get("strike")
    if strike in ("sngStrike", "dblStrike"):
        out["strike"] = True
    baseline = rPr.get("baseline")
    if baseline is not None:
        try:
            b = int(baseline)
            if b < 0:
                out["sub"] = True
            elif b > 0:
                out["sup"] = True
        except ValueError:
            pass
    if rPr.find(qn("a:highlight")) is not None:
        out["highlight"] = True
    return out


def _defrpr_props(defRPr, theme):
    """Extract run props {bold,italic,underline,color,fontFamily,fontSize} from an a:defRPr element.

    Generated decks (and python-pptx-authored ones) frequently put the real styling on the
    PARAGRAPH's a:pPr/a:defRPr (default run properties) while leaving the actual <a:r> empty, so a
    run with no own size/colour must inherit these, or titles import tiny and uncoloured.
    """
    out = {}
    if defRPr is None:
        return out
    if defRPr.get("b") == "1":
        out["bold"] = True
    if defRPr.get("i") == "1":
        out["italic"] = True
    u = defRPr.get("u")
    if u and u != "none":
        out["underline"] = True
    sz = defRPr.get("sz")
    if sz:
        try:
            out["fontSize"] = round(int(sz) / 100.0 * PT_TO_PX)
        except ValueError:
            pass
    fill = defRPr.find(qn("a:solidFill"))
    if fill is not None:
        c = _resolve_color_element(_first_color_el(fill), theme)
        if c:
            out["color"] = c
    latin = defRPr.find(qn("a:latin"))
    if latin is not None and latin.get("typeface"):
        out["fontFamily"] = latin.get("typeface")
    return out


def _paragraph_defaults(para, shape, level, theme):
    """Default run props inherited by a paragraph's runs: its own pPr/defRPr first, then the shape's
    list-style lvl(N) defRPr. Lets empty runs inherit size/colour/bold/font (the common case)."""
    defaults = {}
    # Shape list-style lvl defRPr (lower priority).
    try:
        el = getattr(shape, "_element", None)
        txBody = el.find(qn("p:txBody")) if el is not None else None
        lstStyle = txBody.find(qn("a:lstStyle")) if txBody is not None else None
        if lstStyle is not None:
            lvl = lstStyle.find(qn(f"a:lvl{level + 1}pPr"))
            if lvl is not None:
                defaults.update(_defrpr_props(lvl.find(qn("a:defRPr")), theme))
    except Exception:
        pass
    # Paragraph pPr/defRPr (higher priority, wins over list style).
    try:
        pPr = para._p.find(qn("a:pPr"))
        if pPr is not None:
            defaults.update(_defrpr_props(pPr.find(qn("a:defRPr")), theme))
    except Exception:
        pass
    return defaults


def _run_view(run, theme):
    """One pptx run → our run dict.

    { text, bold, italic, underline, strike, sub, sup, link, highlight, color, fontFamily, fontSize }
    Empty/None values are omitted so the frontend mark-builder only sets what's present.
    """
    font = run.font
    view = {"text": run.text or ""}
    if font.bold:
        view["bold"] = True
    if font.italic:
        view["italic"] = True
    if font.underline:
        view["underline"] = True
    color = _run_color(run, theme)
    if color:
        view["color"] = color
    if font.name:
        view["fontFamily"] = font.name
    if font.size is not None:
        view["fontSize"] = round(font.size.pt * PT_TO_PX)
    view.update(_run_marks(run))
    # Hyperlink (run.hyperlink.address resolves the r:embed relationship to a URL).
    try:
        addr = run.hyperlink.address
        if addr:
            view["link"] = addr
    except Exception:
        pass
    return view


def _paragraph_list(para, default_bullet=False):
    """Detect a paragraph's list kind + level from its pPr (a:buChar / a:buAutoNum).

    Returns (list_kind|None, level int). a:buNone explicitly suppresses a bullet. When
    `default_bullet` is set (the paragraph lives in a body/content placeholder, which bullets by
    default in PowerPoint), a paragraph with NO explicit bullet element AND no a:buNone inherits a
    bullet, this is how most real decks mark up body text (bullets come from the layout/master
    list style, not an explicit a:buChar on every line).
    """
    level = 0
    try:
        level = int(para.level or 0)
    except Exception:
        level = 0
    try:
        pPr = para._p.find(qn("a:pPr"))
    except Exception:
        pPr = None
    if pPr is not None:
        if pPr.find(qn("a:buNone")) is not None:
            return None, level
        if pPr.find(qn("a:buAutoNum")) is not None:
            return "number", level
        if pPr.find(qn("a:buChar")) is not None:
            return "bullet", level
    # No explicit bullet element: body placeholders inherit a bullet; everything else stays plain.
    return ("bullet" if default_bullet else None), level


def _placeholder_role(shape):
    """'title' | 'body' | 'other' | None for a placeholder shape (None = not a placeholder).

    Body/content/subtitle/object placeholders bullet their text by default; titles and plain text
    boxes do not.
    """
    try:
        ph_type = shape.placeholder_format.type
    except Exception:
        return None
    if ph_type is None:
        return None
    name = getattr(ph_type, "name", "") or ""
    if "TITLE" in name:
        return "title"
    if any(k in name for k in ("BODY", "SUBTITLE", "OBJECT", "CONTENT")):
        return "body"
    return "other"


def _effective_placeholder_size(shape, level):
    """Resolve an inherited font size (px) for a placeholder paragraph at `level`.

    Chain (gap #16): the shape's own txBody defRPr → layout placeholder → master placeholder →
    master p:txStyles (titleStyle/bodyStyle/otherStyle) by level. Returns px or None.
    """
    ph_type = None
    try:
        ph_type = shape.placeholder_format.type
    except Exception:
        return None
    if ph_type is None:
        return None

    name = getattr(ph_type, "name", "") or ""
    is_title = "TITLE" in name or name == "CENTER_TITLE"

    # 1) The shape's own list-style defRPr for this level.
    sz = _list_style_size(getattr(shape, "_element", None), level)
    if sz:
        return round(sz * PT_TO_PX)

    # 2) Layout then master placeholder with the matching idx, their list-style defRPr.
    try:
        idx = shape.placeholder_format.idx
    except Exception:
        idx = None
    for source in _inherited_placeholders(shape, idx):
        sz = _list_style_size(getattr(source, "_element", source), level)
        if sz:
            return round(sz * PT_TO_PX)

    # 3) Master p:txStyles by role + level.
    sz = _master_txstyle_size(shape, is_title, name, level)
    if sz:
        return round(sz * PT_TO_PX)
    return None


def _inherited_placeholders(shape, idx):
    """Yield the layout then master placeholder(s) that `shape` inherits from (by idx)."""
    layout = None
    master = None
    try:
        # shape.part is the slide part; reach the slide → layout → master.
        from pptx.parts.slide import SlidePart

        if isinstance(shape.part, SlidePart):
            layout = shape.part.slide.slide_layout
            master = layout.slide_master
    except Exception:
        pass
    for container in (layout, master):
        if container is None:
            continue
        try:
            for ph in container.placeholders:
                if idx is None or ph.placeholder_format.idx == idx:
                    yield ph
                    break
        except Exception:
            continue


def _list_style_size(element, level):
    """Read txBody/lstStyle/lvlNpPr/defRPr@sz (in pt) for `level` off a shape element, or None."""
    if element is None:
        return None
    try:
        txBody = element.find(qn("p:txBody"))
        if txBody is None:
            return None
        lstStyle = txBody.find(qn("a:lstStyle"))
        if lstStyle is None:
            return None
        lvl = lstStyle.find(qn(f"a:lvl{level + 1}pPr"))
        if lvl is None:
            return None
        defRPr = lvl.find(qn("a:defRPr"))
        if defRPr is None:
            return None
        sz = defRPr.get("sz")
        return int(sz) / 100.0 if sz else None
    except Exception:
        return None


def _master_txstyle_size(shape, is_title, ph_name, level):
    """Master p:txStyles size (pt) for the placeholder role + level, or None."""
    try:
        master = shape.part.slide.slide_layout.slide_master
    except Exception:
        return None
    try:
        txStyles = master.element.find(qn("p:txStyles"))
        if txStyles is None:
            return None
        if is_title:
            style = txStyles.find(qn("p:titleStyle"))
        elif (
            "BODY" in (ph_name or "")
            or "SUBTITLE" in (ph_name or "")
            or "OBJECT" in (ph_name or "")
        ):
            style = txStyles.find(qn("p:bodyStyle"))
        else:
            style = txStyles.find(qn("p:otherStyle"))
            if style is None:
                style = txStyles.find(qn("p:bodyStyle"))
        if style is None:
            return None
        lvl = style.find(qn(f"a:lvl{level + 1}pPr"))
        if lvl is None:
            lvl = style.find(qn("a:lvl1pPr"))
        if lvl is None:
            return None
        defRPr = lvl.find(qn("a:defRPr"))
        if defRPr is None:
            return None
        sz = defRPr.get("sz")
        return int(sz) / 100.0 if sz else None
    except Exception:
        return None


def _has_visible_text(tf):
    """True if any run in the text frame has non-whitespace text (gap #5b)."""
    try:
        for para in tf.paragraphs:
            for run in para.runs:
                if (run.text or "").strip():
                    return True
    except Exception:
        return True  # be lenient on error, don't silently drop content
    return False


def _text_object(shape, theme):
    """A shape with a text frame → a {type:'text', ...} view with rich paragraphs.

    The box also carries fallback defaults (align/valign + the first run's font/size/color) so a
    text box still looks right even if individual runs carry no explicit font. Inherited
    placeholder sizes (titles/body) are resolved when runs have no explicit size.
    """
    tf = shape.text_frame
    paragraphs = []
    first_font = None  # remember the first styled run for box-level fallbacks
    box_align = None
    box_line_height = None
    # Body/content placeholders bullet their text by default (inherited from the layout/master),
    # even with no explicit a:buChar, so default paragraphs there to bullets.
    role = _placeholder_role(shape)
    default_bullet = role == "body"
    # Theme latin font a run inherits when it sets no explicit typeface (the +mj-lt/+mn-lt case):
    # titles use the major (heading) font, everything else the minor (body) font. Lets theme-Calibri
    # text resolve to a real name so fontStack() can map it to the Carlito metric clone.
    theme_font = theme.theme_font(role) if theme is not None else None

    # Text-frame properties: insets (padding) + autofit. normAutofit bakes a font shrink + line
    # spacing reduction into the layout, so we scale every effective run size by it to match PowerPoint.
    body_pr = _body_pr(shape)
    insets = _text_insets(body_pr)
    autofit = _autofit_kind(body_pr)
    # bodyPr wrap="none": PowerPoint does NOT auto-wrap, each paragraph stays on one line and
    # overflows the box. Our renderer wraps to the box width by default, so without this every line
    # break differs from the ground truth. Default (absent / "square") wraps normally.
    nowrap = body_pr is not None and body_pr.get("wrap") == "none"
    font_scale, ln_reduction = (
        _norm_autofit_scale(body_pr) if autofit == "norm" else (1.0, 0.0)
    )

    def _scaled(px):
        return round(px * font_scale) if (font_scale != 1.0 and px) else px

    for para in tf.paragraphs:
        align = _ALIGN.get(para.alignment)
        if box_align is None and align:
            box_align = align

        spacing = _para_spacing(para)
        # Box-level lineHeight fallback: the first percent line-spacing we see (or python-pptx's float
        # multiple). Per-paragraph spacing on the paragraph view wins; this is the box default.
        if box_line_height is None and "lineHeight" in spacing:
            box_line_height = spacing["lineHeight"]
        if box_line_height is None:
            ls = para.line_spacing
            # Only a float multiple (e.g. 1.5) maps to a unitless lineHeight; Pt/Length is skipped.
            if isinstance(ls, (int, float)) and not hasattr(ls, "pt"):
                box_line_height = float(ls)

        list_kind, level = _paragraph_list(para, default_bullet=default_bullet)

        runs = [_run_view(r, theme) for r in para.runs]
        # A paragraph with no runs but with text (rare) still yields an empty line.
        if not runs:
            runs = [{"text": ""}]
        # Don't bullet a blank line (an inherited default bullet only makes sense with text).
        if list_kind and not any((r.get("text") or "").strip() for r in runs):
            list_kind = None

        # Inherit paragraph-default run props (a:pPr/a:defRPr or shape list style) for runs that
        # don't set their own, the title/body styling in generated decks lives here, not on <a:r>.
        pdefaults = _paragraph_defaults(para, shape, level, theme)
        if pdefaults:
            for r in runs:
                for k, v in pdefaults.items():
                    r.setdefault(k, v)

        # Theme-font fallback: a run with NO explicit font (none on the run, none in defRPr/list-style)
        # inherits the theme major/minor latin typeface. Resolve it so the renderer & ground truth use
        # the same metric-compatible clone instead of each substituting a different fallback.
        if theme_font:
            for r in runs:
                r.setdefault("fontFamily", theme_font)

        # Resolve inherited size for runs that carry none (placeholders: title/body).
        eff_size = None
        if any("fontSize" not in r for r in runs):
            eff_size = _effective_placeholder_size(shape, level)
        if eff_size:
            for r in runs:
                r.setdefault("fontSize", eff_size)

        # normAutofit shrink: scale every run's effective size so text fits the authored box exactly
        # as PowerPoint renders it (font is NOT mapped here, only the size is scaled).
        if font_scale != 1.0:
            for r in runs:
                if r.get("fontSize"):
                    r["fontSize"] = _scaled(r["fontSize"])

        if first_font is None:
            for r in runs:
                if r.get("fontFamily") or r.get("fontSize") or r.get("color"):
                    first_font = r
                    break

        para_view = {"heading": 0, "align": align, "runs": runs}
        if list_kind:
            para_view["list"] = list_kind
            para_view["level"] = level
        elif level:
            para_view["level"] = level
        # Per-paragraph spacing (lnSpc / spcBef / spcAft). lnSpcReduction from normAutofit further
        # tightens the line box, matching PowerPoint's autofit.
        if "lineHeight" in spacing:
            lh = spacing["lineHeight"] - ln_reduction
            para_view["lineHeight"] = round(max(0.1, lh), 4)
        if "lineHeightPx" in spacing:
            para_view["lineHeightPx"] = spacing["lineHeightPx"]
        if "spaceBefore" in spacing:
            para_view["spaceBefore"] = spacing["spaceBefore"]
        if "spaceAfter" in spacing:
            para_view["spaceAfter"] = spacing["spaceAfter"]
        paragraphs.append(para_view)

    view = {
        "type": "text",
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "w": emu_to_px(shape.width),
        "h": emu_to_px(shape.height),
        "rotation": float(shape.rotation or 0),
        "valign": _vanchor(tf),
        "paragraphs": paragraphs,
        # Text-box padding from bodyPr insets (PowerPoint defaults applied when unset). This shifts
        # where text starts and changes the wrap width, both needed to match the ground truth.
        "padTop": insets["top"],
        "padRight": insets["right"],
        "padBottom": insets["bottom"],
        "padLeft": insets["left"],
    }
    # Autofit mode: 'shape' (spAutoFit), authored box IS the fitted size, so don't clip; the box
    # grows to its content. 'norm' is already baked into the scaled run sizes above. The renderer
    # reads this to decide overflow behaviour without clipping text PowerPoint shows in full.
    if autofit:
        view["autofit"] = autofit
    if nowrap:
        view["nowrap"] = True
    if box_align:
        view["align"] = box_align
    # Box-level spacing: when every paragraph shares the same spacing value, lift it to the box so
    # the renderer applies it uniformly via CSS (robust against the rich-text round-trip). Per-paragraph
    # values still live on the paragraphs for export / future non-uniform decks.
    for key in ("lineHeight", "lineHeightPx", "spaceBefore", "spaceAfter"):
        vals = {p[key] for p in paragraphs if key in p}
        if len(vals) == 1:
            view[key] = next(iter(vals))
    if "lineHeight" not in view and box_line_height:
        view["lineHeight"] = box_line_height
    # Default single line spacing: PowerPoint's "single" is the font's natural line box (~1.2 for the
    # body fonts here), NOT the CSS/browser default (~1.5). Pin imported boxes that set no explicit
    # lnSpc to this so multi-line bodies don't drift downward vs the ground truth. (Editor-created
    # boxes never hit this path, it's import-only.)
    if "lineHeight" not in view and "lineHeightPx" not in view:
        view["lineHeight"] = _SINGLE_LINE_FACTOR
    # Box-level font/color fallbacks (the renderer reads these when a run has no own mark).
    if first_font:
        if first_font.get("fontFamily"):
            view["fontFamily"] = first_font["fontFamily"]
        if first_font.get("fontSize"):
            view["fontSize"] = first_font["fontSize"]
        if first_font.get("color"):
            view["fill"] = first_font["color"]
    return view


# ---------------------------------------------------------------------------
# Shape fill / stroke / effects (gaps #5a, #8, #9, #10, #11, #12)
# ---------------------------------------------------------------------------


def _sp_pr(shape):
    """The shape's <p:spPr> element, or None."""
    try:
        spPr = shape._element.find(qn("p:spPr"))
        if spPr is None:
            spPr = shape._element.find(qn("a:spPr"))
        return spPr
    except Exception:
        return None


def _rgba_or_hex(color_el, theme):
    """A color element → an "rgba(...)" string when it has alpha, else its "#rrggbb"."""
    hex_color = _resolve_color_element(color_el, theme)
    if not hex_color:
        return None
    a = _color_alpha(color_el)
    if a >= 1.0:
        return hex_color
    r, g, b = _hex_to_rgb(hex_color)
    return f"rgba({r},{g},{b},{a})"


def _gradient_css(grad_fill, theme):
    """An <a:gradFill> → a CSS gradient string from its stops, or None.

    Handles a linear gradient (<a:lin ang>) and a radial/path gradient (<a:path>, emitted as a CSS
    radial-gradient). Per-stop alpha is preserved as rgba(). Returns None for <2 usable stops.
    """
    try:
        gs_lst = grad_fill.find(qn("a:gsLst"))
        if gs_lst is None:
            return None
        stops = []
        for gs in gs_lst.findall(qn("a:gs")):
            pos = gs.get("pos")
            frac = (int(pos) / 1000.0) if pos is not None else None  # pos is /1000 %
            color = _rgba_or_hex(_first_color_el(gs), theme)
            if color:
                stops.append((frac, color))
        if len(stops) < 2:
            return None
        parts = [f"{c} {round(f, 1)}%" if f is not None else c for f, c in stops]
        body = ", ".join(parts)
        # Radial/path gradient → CSS radial-gradient (centre-out is the common authoring case).
        if grad_fill.find(qn("a:path")) is not None:
            return f"radial-gradient(circle, {body})"
        # Linear gradient: <a:lin ang="..."> (1/60000 degree, clockwise from 3 o'clock). CSS 0deg = to
        # top and grows clockwise, OOXML 0 = to the right, so CSS angle = ooxml + 90. Default top→bottom.
        angle_deg = 90
        lin = grad_fill.find(qn("a:lin"))
        if lin is not None and lin.get("ang") is not None:
            ooxml_ang = int(lin.get("ang")) / 60000.0
            angle_deg = round((ooxml_ang + 90) % 360)
        return f"linear-gradient({angle_deg}deg, {body})"
    except Exception:
        return None


def _gradient_spec(grad_fill, theme):
    """An <a:gradFill> → structured {kind, angle, stops:[{color,pos}]} for SVG path fills.

    Custom-geometry (SVG <path>) shapes can't take a CSS gradient string, so they carry this and
    the renderer builds an SVG linear/radial gradient. None for <2 usable stops.
    """
    try:
        gs_lst = grad_fill.find(qn("a:gsLst"))
        if gs_lst is None:
            return None
        stops = []
        for i, gs in enumerate(gs_lst.findall(qn("a:gs"))):
            pos = gs.get("pos")
            frac = (int(pos) / 1000.0) if pos is not None else (i * 100.0)
            color = _rgba_or_hex(_first_color_el(gs), theme)
            if color:
                stops.append({"color": color, "pos": round(frac, 1)})
        if len(stops) < 2:
            return None
        if grad_fill.find(qn("a:path")) is not None:
            return {"kind": "radial", "angle": 0, "stops": stops}
        angle = 90
        lin = grad_fill.find(qn("a:lin"))
        if lin is not None and lin.get("ang") is not None:
            angle = round((int(lin.get("ang")) / 60000.0 + 90) % 360)
        return {"kind": "linear", "angle": angle, "stops": stops}
    except Exception:
        return None


def _solid_alpha(solid_fill):
    """A solidFill's color alpha as 0..1 opacity, or None if fully opaque/unspecified."""
    try:
        color_el = _first_color_el(solid_fill)
        if color_el is None:
            return None
        alpha = color_el.find(qn("a:alpha"))
        if alpha is not None and alpha.get("val") is not None:
            return round(int(alpha.get("val")) / 100000.0, 3)  # val is /1000 %
    except Exception:
        pass
    return None


def _color_alpha(color_el):
    """A color element's <a:alpha> as 0..1, or 1.0 (fully opaque) when unset."""
    if color_el is None:
        return 1.0
    try:
        alpha = color_el.find(qn("a:alpha"))
        if alpha is not None and alpha.get("val") is not None:
            return round(int(alpha.get("val")) / 100000.0, 3)
    except Exception:
        pass
    return 1.0


def _shadow_css(effect_lst, theme):
    """An <a:effectLst> with an <a:outerShdw> → a CSS drop-shadow string, or None.

    Reads the real OOXML shadow geometry so shadowed shapes match PowerPoint instead of a fixed
    drop-shadow:
      - blurRad / dist : EMU → px
      - dir            : 1/60000 degree, clockwise from 3 o'clock → x/y offset
      - color + alpha  : the shadow colour (default black) at its alpha (default ~opaque)
    Returns e.g. "drop-shadow(-5px 0px 11px rgba(0,0,0,0.30))". None when there's no outer shadow.
    """
    if effect_lst is None:
        return None
    shdw = effect_lst.find(qn("a:outerShdw"))
    if shdw is None:
        return None
    try:
        blur_emu = int(shdw.get("blurRad") or 0)
        dist_emu = int(shdw.get("dist") or 0)
        dir_60k = int(shdw.get("dir") or 0)
    except ValueError:
        return None
    blur_px = round(blur_emu / EMU_PER_PX)
    dist_px = dist_emu / EMU_PER_PX
    angle = math.radians(
        dir_60k / 60000.0
    )  # clockwise from east; screen y grows downward
    ox = round(dist_px * math.cos(angle))
    oy = round(dist_px * math.sin(angle))

    color_el = _first_color_el(shdw)
    hex_color = _resolve_color_element(color_el, theme) or "#000000"
    r, g, b = _hex_to_rgb(hex_color)
    a = _color_alpha(color_el)
    return f"drop-shadow({ox}px {oy}px {blur_px}px rgba({r},{g},{b},{a}))"


def _style_ref_color(shape, theme, ref_tag):
    """Colour from a shape's <p:style>/<a:fillRef|a:lnRef> (themed fill/line). None if absent.

    Many themed shapes carry no direct solidFill, their colour lives in the style matrix
    reference. idx="0" means "no fill", so we ignore that case.
    """
    try:
        style = shape._element.find(qn("p:style"))
        ref = style.find(qn(ref_tag)) if style is not None else None
        if ref is None or ref.get("idx") in (None, "0"):
            return None
        return _resolve_color_element(_first_color_el(ref), theme)
    except Exception:
        return None


def _shape_fill_stroke(shape, theme):
    """Best-effort fill/stroke for an auto shape, reading theme colors + effects off spPr.

    Returns a dict with any of: fill, gradient, stroke, strokeWidth, opacity, shadow, radius.
    Only sets stroke/strokeWidth when a real line color/width is authored (gap #5a, no spurious
    1px border on shapes with inherited/none outline).
    """
    out = {}
    spPr = _sp_pr(shape)

    # Explicit <a:noFill> → transparent. Emit 'none' so the renderer does NOT fall back to its
    # editor default fill (which would paint a genuinely transparent imported shape pale blue).
    if spPr is not None and spPr.find(qn("a:noFill")) is not None:
        out["fill"] = "none"
    elif spPr is not None:
        # Fill: solid (theme-aware) or gradient.
        grad = spPr.find(qn("a:gradFill"))
        if grad is not None:
            css = _gradient_css(grad, theme)
            if css:
                out["gradient"] = css
        solid = spPr.find(qn("a:solidFill"))
        if solid is not None and "gradient" not in out:
            color = _resolve_color_element(_first_color_el(solid), theme)
            if color:
                out["fill"] = color
            alpha = _solid_alpha(solid)
            if alpha is not None and alpha < 1.0:
                out["opacity"] = alpha

    # Fall back to python-pptx solid fill, then the themed style-matrix fillRef, if nothing yet.
    if "fill" not in out and "gradient" not in out:
        try:
            if shape.fill.type is not None and shape.fill.type == 1:  # MSO_FILL.SOLID
                fb = _hex(shape.fill.fore_color)
                if fb:
                    out["fill"] = fb
        except Exception:
            pass
    if "fill" not in out and "gradient" not in out:
        ref = _style_ref_color(shape, theme, "a:fillRef")
        if ref:
            out["fill"] = ref

    # Stroke: only when an explicit line color or width is authored.
    if spPr is not None:
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            no_fill = ln.find(qn("a:noFill"))
            ln_solid = ln.find(qn("a:solidFill"))
            if no_fill is None and ln_solid is not None:
                color = _resolve_color_element(_first_color_el(ln_solid), theme)
                if color:
                    out["stroke"] = color
            w = ln.get("w")
            if w is not None and no_fill is None:
                try:
                    out["strokeWidth"] = max(1, emu_to_px(int(w)))
                except ValueError:
                    pass
            dash = _line_dash(ln)
            if dash and no_fill is None:
                out["dash"] = dash

    # Shadow: parse the real outerShdw geometry into a CSS drop-shadow so shadowed shapes match
    # PowerPoint's subtle offset/blur instead of our fixed boolean drop-shadow. Keep the boolean too
    # so editor-created shapes (no shadowCss) still get the default shadow.
    if spPr is not None:
        effect_lst = spPr.find(qn("a:effectLst"))
        css = _shadow_css(effect_lst, theme)
        if css:
            out["shadowCss"] = css

    return out


def _round_kind_for(shape):
    """True-ish ('roundRect') if the shape's prstGeom is a rounded-rectangle preset, else None.

    Used for a roundRect that ALSO carries text (so it flows through the text path, where
    auto_shape_type is unreliable), we still want its background corner radius.
    """
    try:
        spPr = _sp_pr(shape)
        prstGeom = spPr.find(qn("a:prstGeom")) if spPr is not None else None
        prst = prstGeom.get("prst") if prstGeom is not None else None
    except Exception:
        prst = None
    if prst in ("roundRect", "round1Rect", "round2SameRect", "round2DiagRect"):
        return "roundRect"
    return None


def _round_radius(shape):
    """For a roundRect, approximate the corner radius in px from its avLst adjustment.

    The adj is a fraction (×100000) of the SHORTER side. Best-effort.
    """
    try:
        spPr = _sp_pr(shape)
        if spPr is None:
            return None
        prstGeom = spPr.find(qn("a:prstGeom"))
        if prstGeom is None:
            return None
        avLst = prstGeom.find(qn("a:avLst"))
        frac = 0.0
        if avLst is not None:
            gd = avLst.find(qn("a:gd"))
            if gd is not None and gd.get("fmla"):
                fmla = gd.get("fmla")  # e.g. "val 16667"
                parts = fmla.split()
                if len(parts) == 2 and parts[0] == "val":
                    frac = int(parts[1]) / 100000.0
        if frac <= 0:
            # PowerPoint default roundRect adjustment is ~16.667%.
            frac = 1.0 / 6.0
        shorter = min(emu_to_px(shape.width), emu_to_px(shape.height))
        return round(frac * shorter)
    except Exception:
        return None


# ---- Group transforms (p:grpSp): map a child's local coords into slide space ----
# A group defines a child coordinate space (chOff/chExt) that maps onto the group's own box
# (off/ext). python-pptx reports a grouped shape's geometry in that CHILD space, so without the
# transform every grouped object lands at the wrong position/size. We carry an affine per axis:
#   slide = a + s * child   ->  stored as (ax, sx, ay, sy).
_IDENTITY_XFRM = (0.0, 1.0, 0.0, 1.0)


def _compose_xfrm(t, g):
    """Compose outer transform t with inner transform g (child -> slide via t∘g)."""
    return (t[0] + t[1] * g[0], t[1] * g[1], t[2] + t[3] * g[2], t[3] * g[3])


def _apply_xfrm(view, t):
    """Map a leaf view's box (x/y/w/h, in its group's child space) into slide space."""
    if view is None or t == _IDENTITY_XFRM or "x" not in view:
        return view
    ax, sx, ay, sy = t
    view["x"] = round(ax + sx * view["x"])
    view["y"] = round(ay + sy * view["y"])
    if "w" in view:
        view["w"] = max(1, round(abs(sx) * view["w"]))
    if "h" in view:
        view["h"] = max(1, round(abs(sy) * view["h"]))
    # A scaled group also scales the font/spacing of text it contains. No-op for layout groups
    # (scale ≈ 1); only kicks in when the group genuinely resizes its children.
    scale = (abs(sx) + abs(sy)) / 2
    if abs(scale - 1) > 0.01:
        for k in (
            "fontSize",
            "lineHeightPx",
            "padTop",
            "padRight",
            "padBottom",
            "padLeft",
            "spaceBefore",
            "spaceAfter",
        ):
            if view.get(k):
                view[k] = round(view[k] * scale, 2)
    return view


def _group_local_xfrm(group_shape):
    """A group's child-space -> parent-space affine from grpSpPr/xfrm (chOff/chExt onto off/ext)."""
    try:
        el = group_shape._element
        grp_pr = el.find(qn("p:grpSpPr"))
        xfrm = grp_pr.find(qn("a:xfrm")) if grp_pr is not None else None
        if xfrm is None:
            return _IDENTITY_XFRM
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        choff = xfrm.find(qn("a:chOff"))
        chext = xfrm.find(qn("a:chExt"))
        if off is None or ext is None or choff is None or chext is None:
            return _IDENTITY_XFRM
        ox = int(off.get("x")) / 9525.0
        oy = int(off.get("y")) / 9525.0
        ex = int(ext.get("cx")) / 9525.0
        ey = int(ext.get("cy")) / 9525.0
        cox = int(choff.get("x")) / 9525.0
        coy = int(choff.get("y")) / 9525.0
        cex = int(chext.get("cx")) / 9525.0
        cey = int(chext.get("cy")) / 9525.0
        sx = ex / cex if cex else 1.0
        sy = ey / cey if cey else 1.0
        return (ox - cox * sx, sx, oy - coy * sy, sy)
    except Exception:
        return _IDENTITY_XFRM


# ---- Custom geometry (a:custGeom) -> an SVG path string in path-unit space ----
def _pt_xy(el):
    """An <a:pt> -> (x, y) ints in path units."""
    return int(el.get("x") or 0), int(el.get("y") or 0)


def _arc_to(cmd, cur):
    """<a:arcTo wR hR stAng swAng> -> (svg 'A' segment, new current point)."""
    try:
        wr = float(cmd.get("wR") or 0)
        hr = float(cmd.get("hR") or 0)
        st = math.radians(float(cmd.get("stAng") or 0) / 60000.0)
        sw = math.radians(float(cmd.get("swAng") or 0) / 60000.0)
    except ValueError:
        return None, cur
    if wr == 0 or hr == 0:
        return None, cur
    cx = cur[0] - wr * math.cos(st)
    cy = cur[1] - hr * math.sin(st)
    end = st + sw
    ex = cx + wr * math.cos(end)
    ey = cy + hr * math.sin(end)
    large = 1 if abs(sw) > math.pi else 0
    sweep = 1 if sw > 0 else 0
    return f"A{round(wr)} {round(hr)} 0 {large} {sweep} {round(ex)} {round(ey)}", (
        ex,
        ey,
    )


def _custom_geom(shape):
    """A shape's <a:custGeom> -> (svg_path_d, path_w, path_h), or None.

    Coordinates stay in the path's own unit space; the renderer scales them to the object box
    via a viewBox + preserveAspectRatio=none. Handles move/line/cubic/quad/arc/close verbs.
    """
    spPr = _sp_pr(shape)
    if spPr is None:
        return None
    cust = spPr.find(qn("a:custGeom"))
    if cust is None:
        return None
    path_lst = cust.find(qn("a:pathLst"))
    if path_lst is None:
        return None
    parts = []
    pw = ph = 0
    for path in path_lst.findall(qn("a:path")):
        try:
            pw = max(pw, int(path.get("w") or 0))
            ph = max(ph, int(path.get("h") or 0))
        except ValueError:
            pass
        cur = (0, 0)
        for cmd in path:
            tag = etree.QName(cmd).localname
            if tag == "moveTo":
                pt = cmd.find(qn("a:pt"))
                if pt is not None:
                    cur = _pt_xy(pt)
                    parts.append(f"M{cur[0]} {cur[1]}")
            elif tag == "lnTo":
                pt = cmd.find(qn("a:pt"))
                if pt is not None:
                    cur = _pt_xy(pt)
                    parts.append(f"L{cur[0]} {cur[1]}")
            elif tag == "cubicBezTo":
                pts = [_pt_xy(p) for p in cmd.findall(qn("a:pt"))]
                if len(pts) == 3:
                    cur = pts[2]
                    parts.append("C" + " ".join(f"{x} {y}" for x, y in pts))
            elif tag == "quadBezTo":
                pts = [_pt_xy(p) for p in cmd.findall(qn("a:pt"))]
                if len(pts) == 2:
                    cur = pts[1]
                    parts.append("Q" + " ".join(f"{x} {y}" for x, y in pts))
            elif tag == "arcTo":
                seg, cur = _arc_to(cmd, cur)
                if seg:
                    parts.append(seg)
            elif tag == "close":
                parts.append("Z")
    if not parts or pw <= 0 or ph <= 0:
        return None
    return " ".join(parts), pw, ph


def _first_grad_color(shape, theme):
    """A representative solid colour (first gradient stop) for a gradient-filled custom shape."""
    spPr = _sp_pr(shape)
    if spPr is None:
        return None
    grad = spPr.find(qn("a:gradFill"))
    gs_lst = grad.find(qn("a:gsLst")) if grad is not None else None
    gs = gs_lst.find(qn("a:gs")) if gs_lst is not None else None
    if gs is None:
        return None
    return _resolve_color_element(_first_color_el(gs), theme)


def _shape_object(shape, theme, skipped):
    """An MSO auto-shape → a {type:'shape', ...} view (or a text view if it carries visible text).

    A filled rectangle with text is, pragmatically, a text box: we emit a text object (carrying
    the shape's fill so the colored background survives). A textless shape stays a shape.
    """
    has_text = shape.has_text_frame and _has_visible_text(shape.text_frame)
    style = _shape_fill_stroke(shape, theme)

    if has_text:
        view = _text_object(shape, theme)
        # Keep the shape's solid fill as the text box background; the run color stays the text color.
        if style.get("fill"):
            view["shapeFill"] = style["fill"]
        # A shape-that-contains-text still owns its border/corner/shadow chrome, forward them so the
        # card behind the text matches (the text layout itself stays the text agent's domain).
        if style.get("stroke"):
            view["stroke"] = style["stroke"]
        if style.get("strokeWidth"):
            view["strokeWidth"] = style["strokeWidth"]
        if style.get("dash"):
            view["dash"] = style["dash"]
        if style.get("shadowCss"):
            view["shadowCss"] = style["shadowCss"]
        # roundRect-with-text: carry the corner radius so the box background rounds like the source.
        kind = _round_kind_for(shape)
        if kind:
            radius = _round_radius(shape)
            if radius:
                view["radius"] = radius
        return view

    # Textless custom geometry (freeform / vector art) → a real SVG path, not a rect fallback.
    cg = _custom_geom(shape)
    if cg is not None:
        d, pw, ph = cg
        view = {
            "type": "shape",
            "shape": "custom",
            "customPath": d,
            "pathW": pw,
            "pathH": ph,
            "x": emu_to_px(shape.left),
            "y": emu_to_px(shape.top),
            "w": emu_to_px(shape.width),
            "h": emu_to_px(shape.height),
            "rotation": float(shape.rotation or 0),
        }
        fill = style.get("fill") or _first_grad_color(shape, theme)
        if fill:
            view["fill"] = fill
        # Real gradient (as SVG stops) so translucent gradient blobs/glows don't flatten to a solid.
        spPr = _sp_pr(shape)
        grad = spPr.find(qn("a:gradFill")) if spPr is not None else None
        if grad is not None:
            spec = _gradient_spec(grad, theme)
            if spec:
                view["gradSpec"] = spec
        for k in ("stroke", "strokeWidth", "dash", "shadowCss"):
            if style.get(k):
                view[k] = style[k]
        if style.get("opacity") is not None:
            view["opacity"] = style["opacity"]
        return view

    try:
        name = shape.auto_shape_type.name if shape.auto_shape_type is not None else None
    except Exception:
        name = None
    kind = _SHAPE_KIND.get(name or "", "rect")
    if name is not None and name not in _SHAPE_KIND:
        skipped["autoshape_fallback_rect"] = (
            skipped.get("autoshape_fallback_rect", 0) + 1
        )

    view = {
        "type": "shape",
        "shape": kind,
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "w": emu_to_px(shape.width),
        "h": emu_to_px(shape.height),
        "rotation": float(shape.rotation or 0),
    }
    if style.get("fill"):
        view["fill"] = style["fill"]
    if style.get("gradient"):
        view["gradient"] = style["gradient"]
    if style.get("stroke"):
        view["stroke"] = style["stroke"]
    if style.get("strokeWidth"):
        view["strokeWidth"] = style["strokeWidth"]
    if style.get("dash"):
        view["dash"] = style["dash"]
    if style.get("opacity") is not None:
        view["opacity"] = style["opacity"]
    if style.get("shadowCss"):
        view["shadowCss"] = style["shadowCss"]
    if kind in _ROUND_KINDS:
        radius = _round_radius(shape)
        if radius:
            view["radius"] = radius
    elif kind == "rect":
        # PowerPoint rectangles have SQUARE corners. The renderer applies an 8px default radius to
        # bare rects (for nicer editor shapes), so emit an explicit radius:0 to keep imported rects
        # crisp without changing editor-created rects (which carry no radius and keep the default).
        view["radius"] = 0
    return view


def _crop_view(shape):
    """A picture's crop (fractions 0..1) → {t,r,b,l} as PERCENT (0..100), or None if uncropped."""
    try:
        crop = {
            "t": round((shape.crop_top or 0) * 100, 2),
            "r": round((shape.crop_right or 0) * 100, 2),
            "b": round((shape.crop_bottom or 0) * 100, 2),
            "l": round((shape.crop_left or 0) * 100, 2),
        }
    except Exception:
        return None
    if any(crop.values()):
        return crop
    return None


def _xfrm_flip(shape):
    """Read a:xfrm flipH/flipV off a shape → (flipH, flipV) booleans."""
    flip_h = flip_v = False
    try:
        spPr = _sp_pr(shape)
        xfrm = spPr.find(qn("a:xfrm")) if spPr is not None else None
        if xfrm is not None:
            flip_h = xfrm.get("flipH") in ("1", "true")
            flip_v = xfrm.get("flipV") in ("1", "true")
    except Exception:
        pass
    return flip_h, flip_v


def _picture_object(shape):
    """A picture → a {type:'image', ...} view with the blob inlined as a base64 data-URL.

    PowerPoint's <a:blipFill> stretches the (srcRect-cropped) image to fill the shape frame,
    ignoring the image's own aspect ratio, so we import with fit:'fill' (the renderer scales the
    cropped region to fill). This matches the LibreOffice ground truth (e.g. a square matplotlib
    PNG dropped into a wide box renders edge-to-edge, not letterboxed).
    """
    image = shape.image
    b64 = base64.b64encode(image.blob).decode()
    content_type = image.content_type or "image/png"
    view = {
        "type": "image",
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "w": emu_to_px(shape.width),
        "h": emu_to_px(shape.height),
        "rotation": float(shape.rotation or 0),
        "src": f"data:{content_type};base64,{b64}",
        "fit": "fill",
    }
    crop = _crop_view(shape)
    if crop:
        view["crop"] = crop
    flip_h, flip_v = _xfrm_flip(shape)
    if flip_h:
        view["flipH"] = True
    if flip_v:
        view["flipV"] = True
    return view


def _table_style(shape, theme):
    """Best-effort table styling: header-row fill, border color, banding flag.

    Reads the header cell fill off the first row (gap #13) and a border color off the table's
    style/cell properties. Returns a dict with any of: fill, stroke, banding.
    """
    out = {}
    try:
        tbl = shape._element.find(".//" + qn("a:tbl"))
    except Exception:
        tbl = None
    if tbl is None:
        return out

    # banding: tblPr@bandRow="1".
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is not None and tblPr.get("bandRow") in ("1", "true"):
        out["banding"] = True

    # Header-row fill: first <a:tr>'s first cell tcPr/solidFill.
    first_tr = tbl.find(qn("a:tr"))
    if first_tr is not None:
        for tc in first_tr.findall(qn("a:tc")):
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is not None:
                solid = tcPr.find(qn("a:solidFill"))
                if solid is not None:
                    color = _resolve_color_element(_first_color_el(solid), theme)
                    if color:
                        out["fill"] = color
                        break

    # Border color: first cell's left/top line solidFill, or any lnL/lnT.
    if first_tr is not None:
        for tc in first_tr.findall(qn("a:tc")):
            tcPr = tc.find(qn("a:tcPr"))
            if tcPr is None:
                continue
            for ln_tag in ("a:lnL", "a:lnT", "a:lnR", "a:lnB"):
                ln = tcPr.find(qn(ln_tag))
                if ln is not None:
                    solid = ln.find(qn("a:solidFill"))
                    if solid is not None:
                        color = _resolve_color_element(_first_color_el(solid), theme)
                        if color:
                            out["stroke"] = color
                            break
            if "stroke" in out:
                break
    return out


def _cell_text_props(tc, theme):
    """Per-cell text alignment / colour / bold, read off the first paragraph + run, or {}.

    Best-effort: tables carry their alignment on the paragraph (a:pPr@algn) and colour on the run
    (a:rPr/a:solidFill). We sample the first non-empty paragraph/run, enough to match the common
    case (a uniformly-styled cell) without a full rich model.
    """
    out = {}
    txBody = tc.find(qn("a:txBody"))
    if txBody is None:
        return out
    for para in txBody.findall(qn("a:p")):
        pPr = para.find(qn("a:pPr"))
        if pPr is not None and pPr.get("algn") and "align" not in out:
            algn = pPr.get("algn")
            out["align"] = {"l": "left", "ctr": "center", "r": "right"}.get(algn)
            if out["align"] is None:
                out.pop("align")
        for run in para.findall(qn("a:r")):
            rPr = run.find(qn("a:rPr"))
            if rPr is None:
                continue
            if "color" not in out:
                solid = rPr.find(qn("a:solidFill"))
                if solid is not None:
                    col = _resolve_color_element(_first_color_el(solid), theme)
                    if col:
                        out["color"] = col
            if "bold" not in out and rPr.get("b") in ("1", "true"):
                out["bold"] = True
            if "color" in out and "bold" in out:
                break
        if "align" in out and "color" in out:
            break
    return out


# OOXML cell vertical anchor (a:tcPr@anchor) → our valign.
_CELL_ANCHOR = {"t": "top", "ctr": "middle", "b": "bottom"}


def _cell_style(tc, theme):
    """One <a:tc> → a CellStyle dict (fill, border, valign, span/merge, text props), or {}."""
    style = {}
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is not None:
        solid = tcPr.find(qn("a:solidFill"))
        if solid is not None:
            col = _resolve_color_element(_first_color_el(solid), theme)
            if col:
                style["fill"] = col
        anchor = tcPr.get("anchor")
        if anchor in _CELL_ANCHOR:
            style["valign"] = _CELL_ANCHOR[anchor]
        # Cell border: take the first defined side that carries a solidFill (uniform-border case).
        for ln_tag in ("a:lnL", "a:lnT", "a:lnR", "a:lnB"):
            ln = tcPr.find(qn(ln_tag))
            if ln is None or ln.find(qn("a:noFill")) is not None:
                continue
            bsolid = ln.find(qn("a:solidFill"))
            if bsolid is not None:
                col = _resolve_color_element(_first_color_el(bsolid), theme)
                if col:
                    style["borderColor"] = col
                    w = ln.get("w")
                    if w is not None:
                        try:
                            style["borderWidth"] = max(1, emu_to_px(int(w)))
                        except ValueError:
                            pass
                    break
    # Merges: gridSpan/rowSpan mark a merge ORIGIN; hMerge/vMerge mark a CONTINUATION (not rendered).
    span = tc.get("gridSpan")
    rowspan = tc.get("rowSpan")
    if span and span.isdigit() and int(span) > 1:
        style["span"] = int(span)
    if rowspan and rowspan.isdigit() and int(rowspan) > 1:
        style["rowSpan"] = int(rowspan)
    if tc.get("hMerge") in ("1", "true") or tc.get("vMerge") in ("1", "true"):
        style["merged"] = True
    style.update(_cell_text_props(tc, theme))
    return style


def _table_geometry(tbl):
    """Relative column widths (a:gridCol@w) and row heights (a:tr@h) in px, or (None, None)."""
    col_widths = []
    grid = tbl.find(qn("a:tblGrid"))
    if grid is not None:
        for gc in grid.findall(qn("a:gridCol")):
            try:
                col_widths.append(emu_to_px(int(gc.get("w"))))
            except (TypeError, ValueError):
                col_widths.append(0)
    row_heights = []
    for tr in tbl.findall(qn("a:tr")):
        try:
            row_heights.append(emu_to_px(int(tr.get("h"))))
        except (TypeError, ValueError):
            row_heights.append(0)
    # Only emit when widths/heights actually vary, equal columns/rows render fine with the default
    # even distribution and don't need to bloat the object.
    cw = col_widths if col_widths and len(set(col_widths)) > 1 else None
    rh = row_heights if row_heights and len(set(row_heights)) > 1 else None
    return cw, rh


def _table_object(shape, theme):
    """A table → a {type:'table', rows, cols, cells, cellStyles, ...} view.

    Captures plain-text cells, table-level styling (header fill / border / banding) AND per-cell
    styling (fill, text colour/align/bold, borders, merges) plus column widths / row heights so
    columns aren't forced equal-width.
    """
    table = shape.table
    rows = list(table.rows)
    cols = list(table.columns)
    cells = []
    for r in range(len(rows)):
        row_cells = []
        for c in range(len(cols)):
            row_cells.append(table.cell(r, c).text)
        cells.append(row_cells)

    view = {
        "type": "table",
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "w": emu_to_px(shape.width),
        "h": emu_to_px(shape.height),
        "rotation": float(shape.rotation or 0),
        "rows": len(rows),
        "cols": len(cols),
        "cells": cells,
    }
    view.update(_table_style(shape, theme))

    # Per-cell styles, indexed [row][col] off the raw <a:tc> elements (so merges/borders survive).
    tbl = shape._element.find(".//" + qn("a:tbl"))
    if tbl is not None:
        cell_styles = []
        any_style = False
        for tr in tbl.findall(qn("a:tr")):
            row_styles = []
            for tc in tr.findall(qn("a:tc")):
                s = _cell_style(tc, theme)
                if s:
                    any_style = True
                row_styles.append(s)
            cell_styles.append(row_styles)
        if any_style:
            view["cellStyles"] = cell_styles
        cw, rh = _table_geometry(tbl)
        if cw:
            view["colWidths"] = cw
        if rh:
            view["rowHeights"] = rh
    return view


# OOXML a:prstDash preset → our coarse dash kind ('dash' | 'dot' | None for solid).
_DASH_KIND = {
    "dot": "dot",
    "sysDot": "dot",
    "dash": "dash",
    "sysDash": "dash",
    "lgDash": "dash",
    "dashDot": "dash",
    "sysDashDot": "dash",
    "lgDashDot": "dash",
    "lgDashDotDot": "dash",
    "sysDashDotDot": "dash",
}


def _line_dash(ln):
    """An <a:ln>'s <a:prstDash> → 'dash' | 'dot' | None (solid/unset). Coarse but renderable."""
    if ln is None:
        return None
    prst = ln.find(qn("a:prstDash"))
    if prst is None:
        return None
    val = prst.get("val")
    if not val or val == "solid":
        return None
    return _DASH_KIND.get(val, "dash")


def _line_arrowheads(shape):
    """Read a connector's head/tail arrowheads off a:ln → (arrowStart, arrowEnd) booleans."""
    arrow_start = False
    arrow_end = False
    try:
        spPr = _sp_pr(shape)
        if spPr is not None:
            ln = spPr.find(qn("a:ln"))
            if ln is not None:
                head = ln.find(qn("a:headEnd"))
                tail = ln.find(qn("a:tailEnd"))
                if head is not None and (head.get("type") or "none") != "none":
                    arrow_start = True
                if tail is not None and (tail.get("type") or "none") != "none":
                    arrow_end = True
    except Exception:
        pass
    return arrow_start, arrow_end


# prstGeom presets that are really straight lines/connectors authored as auto-shapes (so they have
# no text, often zero height), these must render as a stroked line, not an (invisible) rect.
_LINE_PRESETS = {
    "line",
    "straightConnector1",
    "bentConnector2",
    "bentConnector3",
    "curvedConnector3",
}


def _is_line_geom(shape):
    """True if a (non-text) auto-shape's prstGeom is a straight-line/connector preset."""
    try:
        if shape.has_text_frame and _has_visible_text(shape.text_frame):
            return False  # a line never carries body text; keep text shapes on the text path
        spPr = _sp_pr(shape)
        prstGeom = spPr.find(qn("a:prstGeom")) if spPr is not None else None
        prst = prstGeom.get("prst") if prstGeom is not None else None
    except Exception:
        return False
    return prst in _LINE_PRESETS


def _line_object(shape, theme):
    """A connector/line → a {type:'line', ...} view (geometry + stroke + arrowheads)."""
    stroke = None
    stroke_w = None
    spPr = _sp_pr(shape)
    if spPr is not None:
        ln = spPr.find(qn("a:ln"))
        if ln is not None:
            solid = ln.find(qn("a:solidFill"))
            if solid is not None:
                stroke = _resolve_color_element(_first_color_el(solid), theme)
            w = ln.get("w")
            if w is not None:
                try:
                    stroke_w = max(1, emu_to_px(int(w)))
                except ValueError:
                    pass
    if stroke is None:
        try:
            stroke = _hex(shape.line.color)
            if stroke_w is None and shape.line.width is not None:
                stroke_w = max(1, emu_to_px(shape.line.width))
        except Exception:
            pass
    view = {
        "type": "line",
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "w": emu_to_px(shape.width),
        "h": emu_to_px(shape.height),
        "rotation": float(shape.rotation or 0),
    }
    if stroke:
        view["stroke"] = stroke
    if stroke_w:
        view["strokeWidth"] = stroke_w
    if spPr is not None:
        dash = _line_dash(spPr.find(qn("a:ln")))
        if dash:
            view["dash"] = dash
    arrow_start, arrow_end = _line_arrowheads(shape)
    if arrow_start:
        view["arrowStart"] = True
    if arrow_end:
        view["arrowEnd"] = True
    return view


# XL_CHART_TYPE enum-name substrings → our coarse chartType. Anything not matched downgrades to a
# 'column' chart (and is logged), we'd rather show approximate data than drop it.
def _map_chart_type(xl_name):
    """A python-pptx chart_type name → ('bar'|'column'|'line'|'pie'|'area', downgraded?)."""
    n = (xl_name or "").upper()
    if "DOUGHNUT" in n or n.startswith("PIE") or "_PIE" in n:
        # Doughnut renders as a pie in our enum (no doughnut hole support), flag the downgrade.
        return "pie", ("DOUGHNUT" in n)
    if "BAR" in n and "PIE" not in n:
        return "bar", False  # OOXML "bar" = horizontal bars
    if "COLUMN" in n:
        return "column", False
    if "AREA" in n:
        return "area", False
    if "LINE" in n:
        return "line", False
    if "RADAR" in n:
        return "column", True  # no radar renderer, approximate + log
    if "SCATTER" in n or n.startswith("XY"):
        return "line", True  # XY scatter → line (closest), data x dropped to categories
    return "column", True


def _chart_series_color(series, theme):
    """A plot series' fill/line colour as "#rrggbb", or None. Best-effort off the series spPr."""
    try:
        ser = series._element  # CT_BarSer / CT_LineSer / etc.
        spPr = ser.find(qn("c:spPr"))
        if spPr is not None:
            solid = spPr.find(qn("a:solidFill"))
            if solid is not None:
                return _resolve_color_element(_first_color_el(solid), theme)
            ln = spPr.find(qn("a:ln"))
            if ln is not None:
                lsolid = ln.find(qn("a:solidFill"))
                if lsolid is not None:
                    return _resolve_color_element(_first_color_el(lsolid), theme)
    except Exception:
        pass
    return None


def _chart_object(shape, theme, skipped):
    """A graphicFrame holding a c:chart → a {type:'chart', chartType, data, ...} view, or None.

    Reads categories + per-series names/values/colours via python-pptx. Unsupported chart types are
    mapped to the closest of our enum (bar/column/line/pie/area) and the downgrade is recorded in
    `skipped` (so the harness/report shows what was approximated) rather than emitting garbage.
    """
    try:
        chart = shape.chart
    except Exception:
        return None
    try:
        xl_name = getattr(chart.chart_type, "name", str(chart.chart_type))
    except Exception:
        xl_name = None
    chart_type, downgraded = _map_chart_type(xl_name)
    if downgraded:
        key = f"chart_downgrade:{xl_name}->{chart_type}"
        skipped[key] = skipped.get(key, 0) + 1
    # Stacked column/bar: the XL enum name carries it (e.g. COLUMN_STACKED, BAR_STACKED_100).
    stacked = "STACKED" in (xl_name or "").upper()

    # Categories (shared x axis). plots[0] holds them; fall back to indices.
    categories = []
    try:
        plots = list(chart.plots)
        if plots:
            categories = [str(c) if c is not None else "" for c in plots[0].categories]
    except Exception:
        categories = []

    series_out = []
    try:
        for s in chart.series:
            try:
                values = [float(v) if v is not None else 0.0 for v in s.values]
            except Exception:
                values = []
            entry = {
                "name": str(s.name)
                if s.name is not None
                else f"Série {len(series_out) + 1}",
                "values": values,
            }
            col = _chart_series_color(s, theme)
            if col:
                entry["color"] = col
            series_out.append(entry)
    except Exception:
        pass

    if not series_out:
        return None
    if not categories and series_out:
        # No category labels, number the points so the chart still renders.
        categories = [str(i + 1) for i in range(len(series_out[0]["values"]))]

    return {
        "type": "chart",
        "x": emu_to_px(shape.left),
        "y": emu_to_px(shape.top),
        "w": emu_to_px(shape.width),
        "h": emu_to_px(shape.height),
        "rotation": float(shape.rotation or 0),
        "chartType": chart_type,
        "stacked": stacked,
        "data": {"categories": categories, "series": series_out},
    }


def _bg_picture_url(bgPr, part):
    """A <p:bgPr> picture fill (<a:blipFill>/<a:blip r:embed>) → a data-URL, or None.

    Resolves the embedded image relationship against `part` (the slide/layout/master that owns the
    bg XML) and inlines the blob as base64 so the background is fully self-contained.
    """
    try:
        blip_fill = bgPr.find(qn("a:blipFill"))
        if blip_fill is None:
            return None
        blip = blip_fill.find(qn("a:blip"))
        if blip is None:
            return None
        rid = blip.get(qn("r:embed"))
        if not rid or part is None:
            return None
        image_part = part.related_part(rid)
        b64 = base64.b64encode(image_part.blob).decode()
        ct = image_part.content_type or "image/png"
        return f"data:{ct};base64,{b64}"
    except Exception:
        return None


def _bg_from_bgPr(bgPr, theme, part):
    """A <p:bgPr> → a CSS background string (solid hex / gradient / url(...)) or None.

    Returns a value usable directly as the CSS `background` (string-compatible with the existing
    solid-colour contract): a "#rrggbb"/rgba() for solidFill, a "linear/radial-gradient(...)" for
    gradFill, or "url('data:...') center/cover no-repeat" for a picture fill.
    """
    if bgPr is None:
        return None
    grad = bgPr.find(qn("a:gradFill"))
    if grad is not None:
        css = _gradient_css(grad, theme)
        if css:
            return css
    solid = bgPr.find(qn("a:solidFill"))
    if solid is not None:
        color = _rgba_or_hex(_first_color_el(solid), theme)
        if color:
            return color
    url = _bg_picture_url(bgPr, part)
    if url:
        # cover so the picture fills the slide edge-to-edge (PowerPoint stretch/tile → cover is the
        # closest single-declaration match for a full-bleed background).
        return f"url('{url}') center / cover no-repeat"
    return None


def _bg_element(part):
    """The <p:bg> element of a slide/layout/master part, or None."""
    try:
        return part._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    except Exception:
        return None


def _slide_background(slide, theme):
    """Slide background as a CSS string (solid / gradient / picture), resolving inheritance.

    Order: the slide's own <p:bg>, then its layout's, then the master's (a <p:bgRef> style ref on a
    layout/master is left to the theme default and skipped here). Returns a string the frontend can
    drop straight into CSS `background`, or None to inherit the editor default.
    """
    layout = None
    master = None
    try:
        layout = slide.slide_layout
        master = layout.slide_master
    except Exception:
        pass
    for part in (slide, layout, master):
        if part is None:
            continue
        bg = _bg_element(part)
        if bg is None:
            continue
        bgPr = bg.find(qn("p:bgPr"))
        css = _bg_from_bgPr(bgPr, theme, part)
        if css:
            return css
        # A <p:bgRef> on this level points at a theme fill-style we don't expand, stop walking up
        # only if this level explicitly declares noFill; otherwise keep inheriting.
        if bgPr is not None and bgPr.find(qn("a:noFill")) is not None:
            return None

    # Fallback to python-pptx's high-level solid RGB (slide level only).
    try:
        b = slide.background
        if b.fill.type is not None and b.fill.type == 1:  # SOLID
            return _hex(b.fill.fore_color)
    except Exception:
        pass
    return None


def _backdrop_shape_view(shape, part, theme, skipped):
    """A LAYOUT shape → a backdrop view, or None to skip.

    PowerPoint renders the slide layout's graphics behind the slide's own content. We import the
    high-value ones: pictures and picture-filled placeholders (e.g. a full-bleed background image
    that lives on the layout, not the slide) and non-placeholder decorative shapes. Empty text
    placeholders (title/body/footer prompts) are skipped, the slide's own placeholders carry text.
    """
    try:
        st = shape.shape_type
    except Exception:
        st = None
    spPr = _sp_pr(shape)
    # A picture (incl. a picture placeholder, which python-pptx exposes via .image).
    try:
        if st == MSO_SHAPE_TYPE.PICTURE or getattr(shape, "image", None) is not None:
            return _picture_object(shape)
    except Exception:
        pass
    # Text placeholders on the layout are just prompts, skip.
    if getattr(shape, "is_placeholder", False):
        return None
    # Non-placeholder decorative auto-shape / freeform with a visible fill.
    if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        v = _shape_object(shape, theme, skipped)
        has_fill = (
            (v.get("fill") and v.get("fill") != "none")
            or v.get("gradient")
            or v.get("customPath")
        )
        if v and has_fill:
            return v
    return None


def _layout_backdrop(slide, theme, skipped):
    """Backdrop views from the slide's layout (rendered behind the slide's own objects)."""
    out = []
    try:
        layout = slide.slide_layout
    except Exception:
        return out
    for shape in layout.shapes:
        try:
            v = _backdrop_shape_view(shape, layout, theme, skipped)
            if v is not None and not _is_empty_text(v):
                out.append(v)
        except Exception:
            skipped["backdrop_error"] = skipped.get("backdrop_error", 0) + 1
    return out


def _is_empty_text(view):
    """True if a text view has no run with non-whitespace text (gap #5b, skip empties)."""
    if view.get("type") != "text":
        return False
    for para in view.get("paragraphs", []):
        for run in para.get("runs", []):
            if (run.get("text") or "").strip():
                return False
    return True


def _map_shape(shape, objects, theme, skipped, t=_IDENTITY_XFRM):
    """Append the view(s) for one shape to `objects`, or count it as skipped.

    `t` is the affine that maps this shape's local coordinates into slide space (identity for
    top-level shapes; a group's child->slide transform for grouped shapes).
    """
    st = shape.shape_type
    try:
        if st == MSO_SHAPE_TYPE.PICTURE or st == MSO_SHAPE_TYPE.LINKED_PICTURE:
            objects.append(_apply_xfrm(_picture_object(shape), t))
            return
        # Picture PLACEHOLDER (a <p:pic> carrying a <p:ph>), python-pptx types these as PLACEHOLDER,
        # not PICTURE, so they'd otherwise fall through and be dropped. python-pptx still exposes the
        # embedded image via .image (e.g. a full-bleed background photo filling the picture
        # placeholder). Route it through the normal picture path.
        try:
            if getattr(shape, "image", None) is not None:
                objects.append(_apply_xfrm(_picture_object(shape), t))
                return
        except Exception:
            pass
        if getattr(shape, "has_chart", False):
            chart_view = _chart_object(shape, theme, skipped)
            if chart_view is not None:
                objects.append(_apply_xfrm(chart_view, t))
                return
            skipped["chart_empty"] = skipped.get("chart_empty", 0) + 1
            return
        if shape.has_table:
            objects.append(_apply_xfrm(_table_object(shape, theme), t))
            return
        if (
            st == MSO_SHAPE_TYPE.LINE
            or shape.__class__.__name__ == "Connector"
            or _is_line_geom(shape)
        ):
            objects.append(_apply_xfrm(_line_object(shape, theme), t))
            return
        if shape.has_text_frame or st in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
        ):
            # AUTO_SHAPE / FREEFORM (with or without text), TEXT_BOX, PLACEHOLDER flow through here.
            # FREEFORM is custom geometry (a:custGeom), it MUST go to _shape_object or its filled
            # vector path is dropped entirely (e.g. card backgrounds behind white text).
            if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                view = _shape_object(shape, theme, skipped)
            else:
                view = _text_object(shape, theme)
            # Skip empty text boxes (e.g. unfilled layout placeholders), gap #5b.
            if _is_empty_text(view):
                skipped["empty_text"] = skipped.get("empty_text", 0) + 1
                return
            objects.append(_apply_xfrm(view, t))
            return
    except Exception as exc:  # one bad shape must not sink the whole import
        skipped[f"error:{type(exc).__name__}"] = (
            skipped.get(f"error:{type(exc).__name__}", 0) + 1
        )
        return

    # Group shapes, charts, OLE, SmartArt, media, etc., recurse into groups; skip the rest.
    if st == MSO_SHAPE_TYPE.GROUP:
        try:
            ct = _compose_xfrm(t, _group_local_xfrm(shape))
            for child in shape.shapes:
                _map_shape(child, objects, theme, skipped, ct)
            return
        except Exception:
            pass
    key = getattr(st, "name", str(st))
    if st == MSO_SHAPE_TYPE.MEDIA:
        # Movie/audio shapes are not imported (no media extraction); label them by kind so the
        # skipped report says what was dropped instead of a bare "MEDIA".
        try:
            key = "media_audio" if "audioFile" in shape._element.xml else "media_video"
        except Exception:
            key = "media"
    skipped[key] = skipped.get(key, 0) + 1


def parse_pptx(path):
    """Parse a .pptx at `path` into the {mode:'objects', slideSize, slides} contract.

    Raises on a genuinely unreadable file so the caller can fall back to the PNG path.
    """
    prs = Presentation(path)
    slide_size = {
        "w": emu_to_px(prs.slide_width),
        "h": emu_to_px(prs.slide_height),
    }
    skipped = {}
    slides_out = []
    theme_cache = {}  # {id(master): ThemeResolver}, one theme per master.
    for slide in prs.slides:
        theme = ThemeResolver.for_slide(slide, theme_cache)
        objects = []
        # Layout backdrop graphics (e.g. a full-bleed background image that lives on the layout)
        # render behind the slide's own content, so they go first in z-order.
        objects.extend(_layout_backdrop(slide, theme, skipped))
        for shape in slide.shapes:
            _map_shape(shape, objects, theme, skipped)
        slides_out.append(
            {"background": _slide_background(slide, theme), "objects": objects}
        )

    return {
        "mode": "objects",
        "slideSize": slide_size,
        "slides": slides_out,
        "skipped": skipped,  # diagnostic; the frontend ignores it
    }
